"""Generate a 15-minute YouTube deck for the RAG Mastery channel.

Structure:
  - 2-slide RAG channel intro (recreated from the user's draft)
  - 1 architecture slide for Project A (FX Trade Analytics — used as the
    worked example throughout the rest of the video)
  - Heavy on RAG specifics (definition, pipeline, components, hybrid
    search, chunking, evaluation, production patterns, tools)

Run:
    python3 build_rag_intro_deck.py

Output:
    01-RAG-and-FX-Trade-Analytics.pptx (alongside this script)

Source draft for the RAG intro slides (read-only):
    .../01-RAG-Intorduction.pptx (the user's Google Drive)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pathlib import Path

# ---------------------------------------------------------------- design system
NAVY         = RGBColor(0x1E, 0x3A, 0x8A)
NAVY_DARK    = RGBColor(0x0F, 0x2A, 0x5C)
NAVY_SOFT    = RGBColor(0xDB, 0xEA, 0xFE)
INDIGO       = RGBColor(0x31, 0x2E, 0x81)
INDIGO_SOFT  = RGBColor(0xE0, 0xE7, 0xFF)
CYAN         = RGBColor(0x06, 0xB6, 0xD4)
CYAN_SOFT    = RGBColor(0xCF, 0xFA, 0xFE)
EMERALD      = RGBColor(0x10, 0xB9, 0x81)
EMERALD_SOFT = RGBColor(0xD1, 0xFA, 0xE5)
AMBER        = RGBColor(0xF5, 0x9E, 0x0B)
ROSE         = RGBColor(0xF4, 0x3F, 0x5E)

WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
SLATE_50     = RGBColor(0xF8, 0xFA, 0xFC)
SLATE_100    = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_200    = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_400    = RGBColor(0x94, 0xA3, 0xB8)
SLATE_500    = RGBColor(0x64, 0x74, 0x8B)
SLATE_700    = RGBColor(0x33, 0x41, 0x55)
SLATE_900    = RGBColor(0x0F, 0x17, 0x2A)

FONT = 'Inter'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

CHANNEL = 'Training My Hobby · Kishore Veleti'
SERIES  = 'RAG Mastery · Episode 01'

# ---------------------------------------------------------------- helpers
def set_text(tf, text, *, size=18, bold=False, color=SLATE_700, align=PP_ALIGN.LEFT,
             font=FONT, line_spacing=1.15):
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    if not p.runs:
        r = p.add_run()
    else:
        r = p.runs[0]
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, **kwargs):
    box = slide.shapes.add_textbox(x, y, w, h)
    set_text(box.text_frame, text, **kwargs)
    return box


def add_bullets(slide, x, y, w, h, items, *, size=18, color=SLATE_700,
                bullet_color=CYAN, line_spacing=1.4, gap_after=Pt(8)):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0);  tf.margin_bottom = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = gap_after
        rb = p.add_run()
        rb.text = '●  '
        rb.font.name = FONT
        rb.font.size = Pt(size)
        rb.font.color.rgb = bullet_color
        rt = p.add_run()
        rt.text = item
        rt.font.name = FONT
        rt.font.size = Pt(size)
        rt.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, *, fill=WHITE, line=None, line_w=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_rounded_rect(slide, x, y, w, h, *, fill=WHITE, line=None, line_w=Pt(0), radius=0.10):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = radius
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_arrow(slide, x, y, w, h, *, fill=CYAN):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False
    return s


def base_slide(prs, *, page=None, total=None, title=None, eyebrow=None, no_chrome=False):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    if no_chrome:
        return slide
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.10), fill=NAVY)
    add_rect(slide, 0, Inches(0.10), Inches(4.5), Inches(0.04), fill=CYAN)
    add_rect(slide, Inches(4.5), Inches(0.10), Inches(2.0), Inches(0.04), fill=AMBER)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
                CHANNEL, size=10, color=SLATE_400)
    add_textbox(slide, Inches(0.5), Inches(7.20), Inches(8), Inches(0.3),
                SERIES, size=10, color=SLATE_500, bold=True)
    if page is not None:
        add_textbox(slide, Inches(12.0), Inches(7.10), Inches(1.0), Inches(0.3),
                    f'{page:02d} / {total:02d}', size=10, color=SLATE_400, align=PP_ALIGN.RIGHT)
    if eyebrow:
        add_textbox(slide, Inches(0.6), Inches(0.45), Inches(10), Inches(0.4),
                    eyebrow.upper(), size=11, bold=True, color=CYAN)
    if title:
        add_textbox(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.95),
                    title, size=32, bold=True, color=SLATE_900)
        add_rect(slide, Inches(0.6), Inches(1.78), Inches(0.5), Inches(0.06), fill=CYAN)
    return slide


def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ============================================================== slides
def s01_title(prs, _total):
    s = base_slide(prs, no_chrome=True)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY_DARK)
    add_rect(s, 0, 0, Inches(7.5), SLIDE_H, fill=NAVY)
    add_rect(s, 0, Inches(7.30), SLIDE_W, Inches(0.20), fill=CYAN)
    add_rect(s, 0, Inches(7.20), Inches(5.0), Inches(0.10), fill=AMBER)
    add_textbox(s, Inches(0.8), Inches(2.2), Inches(11), Inches(0.4),
                'TRAINING MY HOBBY · YOUTUBE · EPISODE 01', size=14, bold=True, color=CYAN)
    add_textbox(s, Inches(0.8), Inches(2.7), Inches(12), Inches(1.5),
                'RAG Mastery', size=72, bold=True, color=WHITE)
    add_textbox(s, Inches(0.8), Inches(4.1), Inches(12), Inches(0.9),
                'Retrieval-Augmented Generation, end to end',
                size=26, color=CYAN_SOFT)
    add_textbox(s, Inches(0.8), Inches(5.0), Inches(12), Inches(0.6),
                'Worked example: a real-time FX trading platform · ~15 minutes',
                size=18, color=SLATE_200)
    add_textbox(s, Inches(0.8), Inches(6.5), Inches(12), Inches(0.4),
                'Hosted by Kishore Veleti', size=12, color=SLATE_400, bold=True)
    add_speaker_notes(s,
        "Welcome to RAG Mastery, episode one. I'm Kishore, and over the next "
        "fifteen minutes I'm going to give you the clearest possible mental "
        "model for Retrieval-Augmented Generation — what it is, why it exists, "
        "the moving parts you actually need, and how it lands on a real "
        "project. The reason most RAG content feels hand-wavy is that it stops "
        "at the toy notebook. We're not going to do that. We'll use a real-time "
        "FX trading platform that I built end-to-end as the running example, "
        "because real production data has all the messiness — multiple sources, "
        "freshness windows, access controls, exact-match requirements — that "
        "RAG actually has to deal with. By the end of this video you should be "
        "able to explain RAG to a colleague and know exactly what to build "
        "first when you start your own. Let's go.")
    return s


def s02_rag_overview(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='RAG Mastery · channel scope · 1 of 2',
                   title='What this channel covers')
    add_rounded_rect(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.7),
                     fill=NAVY_SOFT, line=None)
    add_textbox(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(0.5),
                'RAG · Retrieval-Augmented Generation', size=22, bold=True, color=NAVY_DARK)
    add_textbox(s, Inches(1.0), Inches(2.95), Inches(11.3), Inches(0.95),
                'Give the LLM the right context at query time — fetched from a knowledge '
                'source you trust. No retraining. Just better, grounded answers.',
                size=15, color=SLATE_700, line_spacing=1.4)

    add_rounded_rect(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.5),
                     fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
    add_textbox(s, Inches(0.95), Inches(4.4), Inches(11.5), Inches(0.4),
                'WHY THIS SERIES EXISTS', size=11, bold=True, color=CYAN)
    items = [
        'Most RAG content is toy examples — we build for real industry data',
        'Same patterns reused across seven verticals (next slide)',
        'End-to-end: ingestion, retrieval, generation, evaluation, ops',
    ]
    for i, t in enumerate(items):
        y = Inches(4.85) + Inches(0.55) * i
        add_textbox(s, Inches(0.95), y, Inches(0.4), Inches(0.4), '●', size=18, color=CYAN)
        add_textbox(s, Inches(1.30), y, Inches(11), Inches(0.5), t, size=15, color=SLATE_900)
    add_speaker_notes(s,
        "Quick channel context before we dive in. If you're new here — RAG "
        "stands for Retrieval-Augmented Generation. The one-line definition: "
        "give the language model the right context at query time, fetched "
        "from a knowledge source you trust. No retraining, no fine-tuning, "
        "no GPUs. Just smarter prompting at runtime. The reason this whole "
        "series exists is that most RAG tutorials online use toy datasets — "
        "the same five Wikipedia articles or a folder of PDFs. That's fine "
        "for a demo, useless for production. We're going to use real industry "
        "data across seven verticals — banking, healthcare, biomed, auto, "
        "education, finance, search — and we'll build the same end-to-end "
        "pipeline for each: ingestion, chunking, embedding, retrieval, "
        "reranking, generation, evaluation, and operations. By the time we're "
        "done with the series you'll have a template you can fork for your "
        "own work.")
    return s


def s03_rag_buckets(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='RAG Mastery · channel scope · 2 of 2',
                   title='Seven verticals · real datasets')
    verticals = [
        ('Auto',       'Driving logs',    INDIGO),
        ('Banking',    'OFAC + AML',      CYAN),
        ('Finance',    'Earnings calls',  EMERALD),
        ('BioMed',     'PubMed corpus',   AMBER),
        ('Healthcare', 'Clinical notes',  ROSE),
        ('Education',  'Textbook + math', NAVY_DARK),
        ('Search',     'Q-passage pairs', CYAN),
    ]
    cols = 4
    card_w = Inches(2.95); card_h = Inches(1.65)
    gap_x = Inches(0.20);   gap_y = Inches(0.20)
    start_x = Inches(0.6);  start_y = Inches(2.2)
    for i, (name, sub, accent) in enumerate(verticals):
        r = i // cols; c = i % cols
        x = start_x + (card_w + gap_x) * c
        y = start_y + (card_h + gap_y) * r
        add_rounded_rect(s, x, y, card_w, card_h, fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
        add_rect(s, x, y, card_w, Inches(0.40), fill=accent)
        add_textbox(s, x + Inches(0.25), y + Inches(0.06), card_w - Inches(0.5), Inches(0.32),
                    f'{i+1:02d}', size=11, bold=True, color=WHITE)
        add_textbox(s, x, y + Inches(0.55), card_w, Inches(0.5),
                    name, size=20, bold=True, color=SLATE_900, align=PP_ALIGN.CENTER)
        add_textbox(s, x, y + Inches(1.05), card_w, Inches(0.4),
                    sub, size=12, color=SLATE_500, align=PP_ALIGN.CENTER)
    add_rounded_rect(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(1.0),
                     fill=SLATE_50, line=SLATE_200, line_w=Pt(0.5))
    add_textbox(s, Inches(0.85), Inches(6.05), Inches(11.5), Inches(0.35),
                'TODAY: BANKING / FINANCE — VIA AN FX TRADING PROJECT',
                size=10, bold=True, color=CYAN)
    add_textbox(s, Inches(0.85), Inches(6.32), Inches(11.5), Inches(0.6),
                'Comment which vertical I should ship next. Code + datasets in the description.',
                size=12, color=SLATE_700, line_spacing=1.4)
    add_speaker_notes(s,
        "Here are the seven verticals we'll cover across the series. Auto, "
        "banking, finance, biomed, healthcare, education, and search. Each "
        "one has its own dataset and its own flavour of retrieval challenge. "
        "Auto is video and time-series logs — wide rows. Banking is sanctions "
        "and AML — exact-match recall matters more than semantic similarity. "
        "Finance is long earnings transcripts where summarization quality "
        "matters. Biomed is dense scientific abstracts. Healthcare is "
        "de-identified clinical notes with strict access controls. Education "
        "is structured textbook content. And search uses MS MARCO-style "
        "question-passage pairs as a baseline benchmark. Today we're in the "
        "banking and finance lane, working with a real FX trading project I "
        "built end-to-end. If a different vertical interests you more, drop "
        "a comment — I prioritize the next episode based on what people ask "
        "for. Now let's actually look at the worked example.")
    return s


def s04_project_arch(prs, page, total):
    """The ONE FX architecture slide. Compact, dense, used as the worked example
    for the rest of the video."""
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Project A · the worked example',
                   title='FX Trade Analytics — architecture')

    boxes = [
        ('Customer\nPortal',    '4201',    CYAN),
        ('Trade\nService',      '8080',    NAVY),
        ('Kafka\ntrade-events', 'Confluent', AMBER),
        ('Risk\nService',       '8081',    NAVY),
        ('Kafka\nenriched',     '',        AMBER),
        ('Indexer',             '8082',    NAVY),
        ('OpenSearch\nfx-trades-{region}', 'Per-region', EMERALD),
    ]
    n = len(boxes)
    box_w = Inches(1.55); arrow_w = Inches(0.20)
    total_w = box_w*n + arrow_w*(n-1)
    start_x = (SLIDE_W - total_w) // 2
    y = Inches(2.4); h = Inches(1.45)
    for i, (label, sub, accent) in enumerate(boxes):
        x = start_x + (box_w + arrow_w) * i
        add_rounded_rect(s, x, y, box_w, h, fill=WHITE, line=accent, line_w=Pt(1.5))
        add_textbox(s, x + Inches(0.05), y + Inches(0.20), box_w - Inches(0.1), Inches(0.7),
                    label, size=11, bold=True, color=SLATE_900, align=PP_ALIGN.CENTER, line_spacing=1.2)
        add_textbox(s, x + Inches(0.05), y + Inches(1.00), box_w - Inches(0.1), Inches(0.4),
                    sub, size=10, color=SLATE_500, align=PP_ALIGN.CENTER)
        if i < n - 1:
            add_arrow(s, x + box_w + Inches(0.01), y + Inches(0.55),
                      arrow_w - Inches(0.02), Inches(0.30), fill=CYAN)

    # Side branches
    add_rounded_rect(s, Inches(0.7), Inches(4.20), Inches(4.0), Inches(0.85),
                     fill=NAVY_SOFT, line=None)
    add_textbox(s, Inches(0.95), Inches(4.30), Inches(3.6), Inches(0.35),
                'MASTER DATA · 8083', size=10, bold=True, color=NAVY_DARK)
    add_textbox(s, Inches(0.95), Inches(4.55), Inches(3.6), Inches(0.45),
                'Postgres + Liquibase · pair allow-list',
                size=11, color=NAVY_DARK)

    add_rounded_rect(s, Inches(4.85), Inches(4.20), Inches(4.0), Inches(0.85),
                     fill=NAVY_SOFT, line=None)
    add_textbox(s, Inches(5.10), Inches(4.30), Inches(3.6), Inches(0.35),
                'ADMIN PORTAL · 4200', size=10, bold=True, color=NAVY_DARK)
    add_textbox(s, Inches(5.10), Inches(4.55), Inches(3.6), Inches(0.45),
                'Currencies · pairs · trade books', size=11, color=NAVY_DARK)

    add_rounded_rect(s, Inches(9.0), Inches(4.20), Inches(3.7), Inches(0.85),
                     fill=NAVY_SOFT, line=None)
    add_textbox(s, Inches(9.25), Inches(4.30), Inches(3.3), Inches(0.35),
                'OBS · 5601 / 3000 / 8085', size=10, bold=True, color=NAVY_DARK)
    add_textbox(s, Inches(9.25), Inches(4.55), Inches(3.3), Inches(0.45),
                'OpenSearch / Grafana / Kafka UI', size=11, color=NAVY_DARK)

    # 4 key facts row
    facts = [
        ('Stack',    'Spring Boot 3.4 · Angular 21'),
        ('Streaming','Confluent Kafka · KRaft mode'),
        ('Search',   'OpenSearch · region-partitioned indices'),
        ('Deploy',   'GitHub Actions + OIDC + CloudFormation'),
    ]
    for i, (k, v) in enumerate(facts):
        x = Inches(0.7) + Inches(3.05) * i
        add_rounded_rect(s, x, Inches(5.4), Inches(2.95), Inches(1.0),
                         fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
        add_textbox(s, x + Inches(0.20), Inches(5.5), Inches(2.55), Inches(0.3),
                    k.upper(), size=10, bold=True, color=CYAN)
        add_textbox(s, x + Inches(0.20), Inches(5.8), Inches(2.55), Inches(0.55),
                    v, size=12, color=SLATE_900, line_spacing=1.3)

    add_speaker_notes(s,
        "This is the entire project on one slide — and that's the only "
        "infrastructure slide you're going to get, because the channel is "
        "about RAG, not about FX trading. Left to right: a customer places "
        "a trade through the customer portal on port 4201; the trade service "
        "on 8080 validates the currency pair against the master-data allow-list; "
        "it drops the validated trade into the Kafka trade-events topic; the "
        "risk service consumes that, computes a risk level — LOW, MEDIUM or "
        "HIGH — and republishes to trade-events-enriched; the indexer consumes "
        "that and writes to OpenSearch — and crucially, partitioned by region, "
        "so we have fx-trades-us-east-1, fx-trades-eu-west-1 and so on. The "
        "side branches: master data on 8083 backed by Postgres and Liquibase; "
        "the admin portal on 4200 for managing reference data; the standard "
        "observability stack — OpenSearch dashboards, Grafana, Kafka UI. "
        "Stack, streaming, search, deploy in the bottom row. From here on, "
        "I'm going to treat this project as the worked example for everything "
        "RAG-related — every concept will land on this concrete pipeline.")
    return s


def s05_bridge(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='From Project A to RAG',
                   title='Why FX trade data is a great RAG worked example')

    cards = [
        ('Multi-source',   'Trades, master data, news, sanctions — exactly the messy mix RAG handles', CYAN),
        ('Real-time',      'Streaming events with strict freshness windows — context goes stale fast', AMBER),
        ('High stakes',    'Wrong answer = financial or compliance impact — citations matter', ROSE),
        ('Multi-tenant',   'Region-partitioned data with access controls — same shape as enterprise RAG', EMERALD),
    ]
    cols = 2
    card_w = Inches(6.0); card_h = Inches(1.95)
    gap = Inches(0.10); start_x = Inches(0.6); start_y = Inches(2.2)
    for i, (head, body, accent) in enumerate(cards):
        r = i // cols; c = i % cols
        x = start_x + (card_w + gap) * c
        y = start_y + (card_h + gap) * r
        add_rounded_rect(s, x, y, card_w, card_h, fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
        add_rect(s, x, y, Inches(0.10), card_h, fill=accent)
        add_textbox(s, x + Inches(0.35), y + Inches(0.25), card_w - Inches(0.5), Inches(0.4),
                    head, size=18, bold=True, color=SLATE_900)
        add_textbox(s, x + Inches(0.35), y + Inches(0.75), card_w - Inches(0.5), Inches(1.1),
                    body, size=13, color=SLATE_700, line_spacing=1.4)
    add_textbox(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.4),
                'These four properties show up in nearly every enterprise RAG problem.',
                size=13, color=SLATE_500, align=PP_ALIGN.CENTER, bold=True)
    add_speaker_notes(s,
        "Why FX trade data as the worked example for RAG? Four reasons "
        "specifically. One — multi-source. A trade isn't just a row; it has "
        "context from master data, sanctions lists, news feeds, internal "
        "policy. RAG was literally invented to weave together signals from "
        "different sources at query time. Two — real-time. Trades stream in "
        "constantly; an answer based on yesterday's data is useless. RAG "
        "with a properly indexed live store handles freshness in a way "
        "fine-tuning never can. Three — high stakes. If the model "
        "hallucinates which trades were high-risk, that's a compliance "
        "incident, not a quirky output. Citations and grounded answers are "
        "non-negotiable. Four — multi-tenant. Region-partitioned data with "
        "access controls is the same shape as enterprise RAG over per-customer "
        "knowledge bases, per-team wikis, per-tenant tickets. These four "
        "properties together describe nearly every enterprise RAG problem "
        "you'll see in the wild.")
    return s


def s06_what_is_rag(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='RAG core',
                   title='What is RAG — in 30 seconds')

    add_rounded_rect(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.5),
                     fill=NAVY_SOFT, line=None)
    add_textbox(s, Inches(1.0), Inches(2.45), Inches(11.3), Inches(0.6),
        'Retrieve relevant context · stuff it into the prompt · let the LLM answer.',
        size=22, color=NAVY_DARK, bold=True)
    add_textbox(s, Inches(1.0), Inches(3.05), Inches(11.3), Inches(0.55),
        'No retraining. No fine-tuning. Just smarter prompting at runtime.',
        size=15, color=SLATE_700)

    # Two columns
    y = Inches(4.0); h = Inches(2.6)
    add_rounded_rect(s, Inches(0.6), y, Inches(5.95), h, fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
    add_rect(s, Inches(0.6), y, Inches(0.10), h, fill=ROSE)
    add_textbox(s, Inches(0.95), y + Inches(0.25), Inches(5.5), Inches(0.4),
                'WITHOUT RAG', size=11, bold=True, color=ROSE)
    add_textbox(s, Inches(0.95), y + Inches(0.65), Inches(5.5), Inches(0.5),
                '“Show me high-risk USD/INR trades from last hour.”',
                size=14, bold=True, color=SLATE_900)
    add_textbox(s, Inches(0.95), y + Inches(1.25), Inches(5.5), Inches(1.2),
                'Model has no idea what trades exist. It guesses, paraphrases, fabricates.',
                size=13, color=SLATE_500, line_spacing=1.4)

    add_rounded_rect(s, Inches(6.85), y, Inches(5.95), h, fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
    add_rect(s, Inches(6.85), y, Inches(0.10), h, fill=EMERALD)
    add_textbox(s, Inches(7.20), y + Inches(0.25), Inches(5.5), Inches(0.4),
                'WITH RAG', size=11, bold=True, color=EMERALD)
    add_textbox(s, Inches(7.20), y + Inches(0.65), Inches(5.5), Inches(0.5),
                '“Show me high-risk USD/INR trades from last hour.”',
                size=14, bold=True, color=SLATE_900)
    add_textbox(s, Inches(7.20), y + Inches(1.25), Inches(5.5), Inches(1.2),
                'Retrieve 12 matching trades from OpenSearch → cite each → grounded answer.',
                size=13, color=SLATE_500, line_spacing=1.4)

    add_speaker_notes(s,
        "Here's RAG distilled to one line: retrieve relevant context, stuff "
        "it into the prompt, let the model answer using that context. No "
        "retraining. No fine-tuning. Just smarter prompting at runtime. "
        "Look at the comparison. Without RAG, the model sees only the user's "
        "question. It has zero knowledge of your trade book — none of those "
        "trades were ever in its training data, and even if they were, "
        "they're stale. So it does what models do when they don't know: "
        "guesses confidently, paraphrases, sometimes invents trade IDs that "
        "don't exist. With RAG, before we call the model, we hit OpenSearch, "
        "pull the twelve actual trades that match the filter, drop them into "
        "the prompt, and ask the model to summarise WITH CITATIONS. Same "
        "question, completely different outcome. The model becomes a reasoning "
        "engine over your data, not a guesser. That's the entire idea — "
        "everything else in this video is plumbing.")
    return s


def s07_pipeline(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='RAG core',
                   title='The pipeline')

    steps = [
        ('Query',    'User question'),
        ('Embed',    'Text → vector'),
        ('Retrieve', 'Top-K from store'),
        ('Augment',  'Build prompt'),
        ('Generate', 'LLM answer'),
        ('Cite',     'Sources back'),
    ]
    n = len(steps)
    box_w = Inches(1.85); arrow_w = Inches(0.25)
    total_w = box_w*n + arrow_w*(n-1)
    start_x = (SLIDE_W - total_w) // 2
    y = Inches(2.8); h = Inches(1.6)
    accents = [NAVY, NAVY, CYAN, CYAN, EMERALD, AMBER]
    for i, (label, desc) in enumerate(steps):
        x = start_x + (box_w + arrow_w) * i
        add_rounded_rect(s, x, y, box_w, h, fill=WHITE, line=accents[i], line_w=Pt(1.5))
        add_textbox(s, x, y + Inches(0.15), box_w, Inches(0.3),
                    f'{i+1:02d}', size=11, bold=True, color=accents[i], align=PP_ALIGN.CENTER)
        add_textbox(s, x, y + Inches(0.45), box_w, Inches(0.45),
                    label, size=18, bold=True, color=SLATE_900, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.1), y + Inches(0.95), box_w - Inches(0.2), Inches(0.55),
                    desc, size=12, color=SLATE_500, align=PP_ALIGN.CENTER)
        if i < n - 1:
            add_arrow(s, x + box_w + Inches(0.02), y + Inches(0.65), arrow_w - Inches(0.04),
                      Inches(0.30), fill=accents[i])

    add_rounded_rect(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.7),
                     fill=SLATE_50, line=SLATE_200, line_w=Pt(0.5))
    add_textbox(s, Inches(0.85), Inches(5.15), Inches(11.5), Inches(0.4),
                'WHERE TIME AND COST ACTUALLY GO', size=11, bold=True, color=CYAN)
    add_textbox(s, Inches(0.85), Inches(5.5), Inches(11.5), Inches(1.1),
                'Steps 2–3 are sub-100ms. Step 5 is the only call that pays for tokens. '
                'Optimise retrieval quality first; smaller good context > larger bad context.',
                size=14, color=SLATE_700, line_spacing=1.5)
    add_speaker_notes(s,
        "Six steps, and you should be able to draw this from memory by the "
        "end of this video. Step one — query, the user asks a question. Step "
        "two — embed, we run that question through an embedding model and get "
        "back a vector, typically 1536 numbers. Step three — retrieve, we use "
        "that vector to do a nearest-neighbour search in our vector store and "
        "get back the top-K most similar known passages. Step four — augment, "
        "we build a prompt that includes those passages plus the original "
        "question. Step five — generate, the LLM produces an answer grounded "
        "in the context we gave it. Step six — cite, we return both the "
        "answer and the source IDs we used so the UI can link back to "
        "evidence. Cost-wise, steps two and three are sub-100 milliseconds "
        "and effectively free. Step five is the only call that pays for "
        "tokens. So the optimization rule that drives everything else: "
        "smaller, better context beats larger, weaker context every time. "
        "We'll come back to this rule throughout the video.")
    return s


def s08_embeddings(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='RAG component · 1 of 5',
                   title='Embeddings — text into vectors')

    items = [
        'Same model on document side AND query side — never mix',
        '1536 dims (OpenAI ada-3) is fine; bigger ≠ better for most tasks',
        'Domain-specific embeddings (BGE, E5) often beat generic on niche corpora',
        'Re-embed if you swap models — vectors aren’t portable',
    ]
    add_bullets(s, Inches(0.6), Inches(2.2), Inches(7.5), Inches(4),
                items, size=18, color=SLATE_900, gap_after=Pt(14))

    # Right callout
    add_rounded_rect(s, Inches(8.4), Inches(2.2), Inches(4.3), Inches(4.4),
                     fill=NAVY_DARK, line=None)
    add_textbox(s, Inches(8.65), Inches(2.4), Inches(3.9), Inches(0.4),
                'FX EXAMPLE', size=11, bold=True, color=CYAN)
    add_textbox(s, Inches(8.65), Inches(2.75), Inches(3.9), Inches(0.6),
                'Embed each trade as:', size=14, color=WHITE, bold=True)
    add_textbox(s, Inches(8.65), Inches(3.25), Inches(3.9), Inches(2.5),
        '"USD→INR · 50000 · risk=HIGH · '
        'us-east-1 · trader=BOOK1 · 2026-05-02"',
        size=12, color=CYAN_SOFT, font='Menlo', line_spacing=1.5)
    add_textbox(s, Inches(8.65), Inches(5.45), Inches(3.9), Inches(0.9),
        'A single dense vector lets you ask:\n“trades like this one.”',
        size=12, color=SLATE_200, line_spacing=1.4)
    add_speaker_notes(s,
        "Component one of five — embeddings. An embedding model takes text "
        "and turns it into a dense vector of numbers, where similar meanings "
        "produce similar vectors. The single most important rule: use the "
        "same embedding model on the document side and the query side. "
        "Always. If you embed your trades with OpenAI's ada-3 and then query "
        "with sentence-transformers, you'll get garbage results — the vector "
        "spaces are simply different. Beyond that — bigger isn't necessarily "
        "better. 1536 dimensions from OpenAI is plenty for most workloads. "
        "Domain-specific open models like BGE or E5 often beat general-purpose "
        "ones on niche corpora — biomedical text, legal contracts, code. "
        "And if you swap embedding models later, you have to re-embed every "
        "single document; the vectors aren't portable. For our FX example "
        "you can see on the right — we embed each trade as a short canonical "
        "string of its key fields, and suddenly 'find me trades like this "
        "one' becomes a single vector lookup. That's the whole magic.")
    return s


def s09_vector_store(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='RAG component · 2 of 5',
                   title='Vector stores — pick the boring one')

    cols = [
        ('pgvector', 'In your existing Postgres', 'Free · transactional · use until you cant', EMERALD),
        ('Pinecone', 'Managed cloud',             'Zero ops · pay per million vectors',         CYAN),
        ('Weaviate / Qdrant', 'Self-host or cloud', 'Hybrid search built-in · open',           AMBER),
        ('FAISS',    'Library, not a DB',         'Notebook / on-disk only · no auth',          SLATE_500),
    ]
    n = len(cols)
    card_w = Inches(2.95); card_h = Inches(3.5)
    gap = Inches(0.12)
    total_w = card_w*n + gap*(n-1)
    start_x = (SLIDE_W - total_w) // 2
    y = Inches(2.4)
    for i, (head, sub, body, accent) in enumerate(cols):
        x = start_x + (card_w + gap) * i
        add_rounded_rect(s, x, y, card_w, card_h, fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
        add_rect(s, x, y, card_w, Inches(0.45), fill=accent)
        add_textbox(s, x, y + Inches(0.10), card_w, Inches(0.32),
                    head.upper(), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.2), y + Inches(0.65), card_w - Inches(0.4), Inches(0.5),
                    sub, size=13, bold=True, color=SLATE_900, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.2), y + Inches(1.25), card_w - Inches(0.4), Inches(2.0),
                    body, size=12, color=SLATE_500, align=PP_ALIGN.CENTER, line_spacing=1.5)

    add_textbox(s, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.5),
        'Default to pgvector. The Postgres you already have is the right answer 80% of the time.',
        size=14, color=NAVY_DARK, align=PP_ALIGN.CENTER, bold=True)
    add_speaker_notes(s,
        "Component two — the vector store. Honest take: start with pgvector "
        "in the Postgres database you already have. It's free, transactional, "
        "indexable, and comfortable up to several million vectors. Most teams "
        "never need more than that. Pinecone and the other managed services "
        "are great when you genuinely don't want to operate a database, when "
        "you need horizontal sharding, or when you're past tens of millions "
        "of vectors. Weaviate and Qdrant are open-source middle-grounds — "
        "self-host or use the cloud, with hybrid search built in, which we'll "
        "discuss in a few slides. FAISS is a library, not a database; great "
        "for notebooks and on-disk inference but no auth, no concurrency, no "
        "ops story. The default rule: pick the boring one. The Postgres you "
        "already run for your application is the right answer eight times "
        "out of ten. Move to managed when scale or ops actually forces you, "
        "not when a blog post tells you to.")
    return s


def s10_retrieval(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='RAG component · 3 of 5',
                   title='Retrieval — top-K, filters, reranking')

    items = [
        ('Top-K vector search',    'Pull the K nearest vectors. K = 20-50 is the usual band.'),
        ('Pre-filter by metadata', 'Region, time window, tenant — apply BEFORE the similarity search'),
        ('Rerank with a cross-encoder', 'Re-score the K candidates with a slower, more accurate model'),
        ('Deduplicate near-clones', 'Otherwise the LLM sees the same fact 5 times'),
    ]
    y = Inches(2.2); h = Inches(0.95)
    for i, (head, body) in enumerate(items):
        ry = y + Inches(1.05) * i
        add_rounded_rect(s, Inches(0.6), ry, Inches(0.95), h, fill=NAVY, line=None)
        add_textbox(s, Inches(0.6), ry + Inches(0.30), Inches(0.95), Inches(0.4),
                    f'{i+1:02d}', size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rounded_rect(s, Inches(1.65), ry, Inches(11.05), h, fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
        add_textbox(s, Inches(1.95), ry + Inches(0.18), Inches(10.5), Inches(0.4),
                    head, size=16, bold=True, color=SLATE_900)
        add_textbox(s, Inches(1.95), ry + Inches(0.55), Inches(10.5), Inches(0.4),
                    body, size=13, color=SLATE_500)
    add_speaker_notes(s,
        "Component three — retrieval. And this is the part that most people "
        "underweight. Retrieval is far more than just vector search. Four "
        "things to layer in. One — top-K vector search itself; pull the K "
        "nearest vectors, where K is typically twenty to fifty, not the "
        "three or five you'll see in tutorials. Two — pre-filter by metadata. "
        "For our FX example, region, time window, trader-book — apply these "
        "filters BEFORE the similarity search, otherwise you waste your K "
        "budget on irrelevant rows. Three — rerank with a cross-encoder. "
        "Pull a wide candidate set with the fast bi-encoder, then re-score "
        "the top fifty with a slower, more accurate cross-encoder model that "
        "can attend to query and document jointly. Massive quality gain for "
        "a small latency cost. Four — deduplicate near-clones. Otherwise the "
        "LLM sees the same fact paraphrased five times and weights it five "
        "times. The single biggest lesson I've learned: retriever quality "
        "matters more than which LLM you pick. A great retriever with GPT-3.5 "
        "beats a weak one with GPT-4. Spend your time here.")
    return s


def s11_augment(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='RAG component · 4 of 5',
                   title='Augment — how you build the prompt')

    add_rounded_rect(s, Inches(0.6), Inches(2.2), Inches(7.6), Inches(4.5),
                     fill=SLATE_900, line=None)
    add_textbox(s, Inches(0.85), Inches(2.35), Inches(7.0), Inches(0.4),
                'Prompt template', size=11, bold=True, color=CYAN, font='Menlo')
    template_lines = [
        'You answer questions about FX trades.',
        'Use ONLY the context below; cite ids.',
        'If the context is empty, say so.',
        '',
        '<context>',
        '{retrieved_trade_1}',
        '{retrieved_trade_2}',
        '...',
        '</context>',
        '',
        'Question: {user_question}',
    ]
    for i, line in enumerate(template_lines):
        add_textbox(s, Inches(0.85), Inches(2.7) + Inches(0.32)*i,
                    Inches(7.2), Inches(0.30),
                    line if line else ' ', size=12, color=SLATE_200, font='Menlo')

    add_rounded_rect(s, Inches(8.5), Inches(2.2), Inches(4.2), Inches(4.5),
                     fill=NAVY_SOFT, line=None)
    add_textbox(s, Inches(8.75), Inches(2.4), Inches(3.8), Inches(0.4),
                'RULES THAT WORK', size=11, bold=True, color=NAVY_DARK)
    add_bullets(s, Inches(8.75), Inches(2.85), Inches(3.8), Inches(3.6),
                ['Cite source IDs',
                 'Cap context tokens',
                 'Order best→worst',
                 'Demand “I don’t know”',
                 'Same template every call'],
                size=13, color=NAVY_DARK, bullet_color=NAVY_DARK, gap_after=Pt(8))
    add_speaker_notes(s,
        "Component four — augmentation. This is just prompt engineering with "
        "discipline. Look at the template on the left. Three things stand "
        "out. First — a strict instruction to use ONLY the context provided. "
        "Without that, the model will quietly fall back to its training data "
        "and you'll have no idea. Second — explicit instruction to cite "
        "source IDs. Citations turn answers into something you can audit. "
        "Third — explicit permission to say 'I don't know' when the context "
        "is empty. Models hate to admit ignorance unless you give them a "
        "graceful exit. Five rules to follow on the right. Cite source IDs. "
        "Cap your context tokens — runaway context costs real money and "
        "actually hurts answer quality past a point. Order retrieved passages "
        "best to worst — recency bias matters. Demand the I-don't-know "
        "answer. And use the same template every call so when you A/B test, "
        "you change one variable at a time. Augmentation isn't glamorous, "
        "but discipline here is what separates a demo from a product.")
    return s


def s12_generate(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='RAG component · 5 of 5',
                   title='Generate — grounded answers + citations')

    items = [
        'Pick a model that supports JSON / structured output for citations',
        'Lower temperature (0–0.3) for factual RAG; higher only for synthesis',
        'Stream the answer; readers tolerate latency they can see',
        'Always return the source IDs used — let the UI link back to evidence',
        'Log the full prompt and the answer — debugging RAG without this is hopeless',
    ]
    add_bullets(s, Inches(0.6), Inches(2.2), Inches(12), Inches(4.5),
                items, size=18, color=SLATE_900, gap_after=Pt(14))
    add_speaker_notes(s,
        "Component five — generation. The fun part, but also the most "
        "expensive. Five practical rules. One — pick a model that supports "
        "structured output, ideally JSON-schema-constrained generation. You "
        "want answers AND citations in machine-readable form, not "
        "free-text you have to regex-parse. Two — temperature low. Zero to "
        "0.3 for factual RAG; higher only when you're explicitly synthesising "
        "or generating creative text. Higher temperature is the enemy of "
        "groundedness. Three — stream the answer to the user. Readers "
        "tolerate a lot more latency when they can see progress. Four — "
        "always emit the source IDs you used. The UI should make these "
        "clickable so users can verify the underlying evidence. Five — log "
        "EVERYTHING. The full prompt, the retrieved IDs, the answer, the "
        "latency, the cost. Debugging RAG without logs is genuinely hopeless. "
        "When something goes wrong in production, your logs are the only "
        "thing that lets you trace it backwards through retrieval, filtering, "
        "and the prompt.")
    return s


def s13_rag_over_fx(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Worked example',
                   title='RAG over FX trades — concrete walkthrough')

    # Query block
    add_rounded_rect(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.0),
                     fill=NAVY_DARK, line=None)
    add_textbox(s, Inches(0.85), Inches(2.35), Inches(11.5), Inches(0.35),
                'USER QUERY', size=10, bold=True, color=CYAN)
    add_textbox(s, Inches(0.85), Inches(2.65), Inches(11.5), Inches(0.45),
                '“Why are we seeing more high-risk USD trades from us-east-1 since Tuesday?”',
                size=16, color=WHITE, font='Menlo')

    # 4-step pipeline applied
    steps = [
        ('Filter',   'WHERE region=us-east-1 AND riskLevel=HIGH AND ts > 2026-04-30'),
        ('Embed',    'Embed the question; nearest 50 trade summaries'),
        ('Rerank',   'Cross-encoder picks the 8 most question-relevant'),
        ('Generate', 'LLM summarises trends + cites trade IDs T-101, T-117, …'),
    ]
    y = Inches(3.55); h = Inches(0.65); gap = Inches(0.08)
    for i, (head, body) in enumerate(steps):
        ry = y + (h + gap) * i
        add_rounded_rect(s, Inches(0.6), ry, Inches(2.0), h, fill=CYAN_SOFT, line=None)
        add_textbox(s, Inches(0.6), ry + Inches(0.18), Inches(2.0), Inches(0.4),
                    head.upper(), size=12, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)
        add_rounded_rect(s, Inches(2.75), ry, Inches(9.95), h, fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
        add_textbox(s, Inches(3.00), ry + Inches(0.18), Inches(9.5), Inches(0.4),
                    body, size=13, color=SLATE_900, font='Menlo' if i == 0 else FONT)

    add_textbox(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4),
        'OpenSearch is both retriever AND data plane — no separate vector DB needed.',
        size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "Now let's put all five components together with our worked example. "
        "User asks: why are we seeing more high-risk USD trades from us-east-1 "
        "since Tuesday? Step one — filter. Apply hard filters first: region "
        "equals us-east-1, riskLevel equals HIGH, timestamp greater than "
        "April thirtieth. This is a basic OpenSearch query and it shrinks "
        "the candidate pool from millions to maybe a few thousand. Step two — "
        "embed. We embed the user's question and find the fifty trades whose "
        "embeddings are closest. Step three — rerank. A cross-encoder picks "
        "the eight most QUESTION-relevant trades from those fifty. Step "
        "four — generate. The LLM gets those eight trades plus the question, "
        "summarises the trend in plain English, and cites trade IDs T-101, "
        "T-117, and so on. The user can click any of those IDs and jump to "
        "the underlying record. The really nice property here is that "
        "OpenSearch is both the retriever AND the data plane. We don't need "
        "a separate vector database — pgvector or Pinecone — because "
        "OpenSearch supports both keyword filtering and dense vector search "
        "natively. One system, one query, one set of access controls.")
    return s


def s14_hybrid(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Technique',
                   title='Hybrid search — keyword + vector')

    add_rounded_rect(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.4),
                     fill=NAVY_SOFT, line=None)
    add_textbox(s, Inches(0.95), Inches(2.4), Inches(11.5), Inches(0.4),
                'WHY YOU NEED BOTH', size=11, bold=True, color=NAVY_DARK)
    add_textbox(s, Inches(0.95), Inches(2.75), Inches(11.5), Inches(0.85),
        'Vector search misses exact identifiers (trade IDs, sanctions list names, ticker symbols). '
        'Keyword search misses paraphrase. Combine them.',
        size=14, color=NAVY_DARK, line_spacing=1.5)

    cols = [
        ('Vector (BM25 misses these)', 'Paraphrase · semantic similarity · "trades like this one"', CYAN),
        ('Keyword (vector misses these)', 'Exact matches · IDs · OFAC names · ticker symbols', AMBER),
        ('Reciprocal Rank Fusion', 'Combine ranks · simple, robust, no hyperparameter tuning', EMERALD),
    ]
    y = Inches(3.95); card_w = Inches(4.0); gap = Inches(0.10); h = Inches(2.7)
    for i, (head, body, accent) in enumerate(cols):
        x = Inches(0.6) + (card_w + gap) * i
        add_rounded_rect(s, x, y, card_w, h, fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
        add_rect(s, x, y, card_w, Inches(0.42), fill=accent)
        add_textbox(s, x, y + Inches(0.08), card_w, Inches(0.32),
                    head.upper(), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.3), y + Inches(0.65), card_w - Inches(0.6), Inches(2.0),
                    body, size=13, color=SLATE_700, line_spacing=1.5)
    add_speaker_notes(s,
        "Hybrid search. This is where a lot of teams get burned. Pure vector "
        "search has a critical weakness — it misses exact identifiers. "
        "Trade ID T-12345, sanctions list name 'OFAC SDN', ticker symbol "
        "AAPL. The vector for 'T-12345' is essentially noise; semantically "
        "meaningless. Pure keyword search has the opposite weakness — it "
        "misses paraphrase. The user types 'large dollar trades' and your "
        "BM25 index has nothing matching the literal word 'large'. The "
        "answer, almost always, is to run both. Vector for paraphrase and "
        "semantic similarity. Keyword for exact matches and proper nouns. "
        "Then combine the two ranked lists. The simplest, most robust way "
        "to combine is Reciprocal Rank Fusion — RRF — which just sums "
        "1/(k+rank) across the two lists. No hyperparameters, no tuning, "
        "and it almost always beats either method alone. If you remember "
        "one technique from this video, make it RRF.")
    return s


def s15_chunking(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Pitfall',
                   title='Chunking — the most underrated decision')

    items = [
        'Chunk size sets retrieval granularity — fixed 512 tokens is rarely right',
        'Overlap (10–15%) helps when answers straddle boundaries',
        'Semantic chunking (split on headings, paragraphs, function defs) > naïve',
        'Store the WHOLE document with each chunk — context for the LLM, not the retriever',
        'For our FX case: one chunk per trade · one chunk per book · one per region',
    ]
    add_bullets(s, Inches(0.6), Inches(2.2), Inches(12), Inches(4.5),
                items, size=17, color=SLATE_900, gap_after=Pt(14))
    add_speaker_notes(s,
        "Chunking. This is the most underrated decision in RAG, full stop. "
        "The default in every framework is something like 'split into 512-token "
        "chunks with 50-token overlap.' That's almost never right for your "
        "actual data. Chunk size sets your retrieval granularity — too small "
        "and the retriever returns fragments without context; too large and "
        "you waste prompt tokens on irrelevance. Overlap of 10-15 percent "
        "helps when answers happen to straddle a boundary. Semantic chunking "
        "— splitting on headings, paragraph breaks, function definitions, "
        "transcript turns — almost always beats naïve fixed-size chunking. "
        "And here's the trick most people miss: store the WHOLE document "
        "alongside each chunk, as metadata. The chunk drives retrieval; the "
        "full document gets passed to the LLM for context. For our FX example "
        "chunking is mercifully simple — one chunk per trade is the natural "
        "unit. Each trade is small, structured, self-contained. Compare that "
        "to chunking a thousand-page legal contract or a four-hour earnings "
        "call transcript and you'll see why this slide exists.")
    return s


def s16_eval(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Quality',
                   title='Evaluation — measure or guess')

    metrics = [
        ('Retrieval', 'Precision@K · Recall@K · MRR', NAVY),
        ('Faithfulness', 'Does the answer match the context?', CYAN),
        ('Answer relevance', 'Does the answer address the question?', EMERALD),
        ('Latency / cost', 'p95 ms · tokens per query · $/100 queries', AMBER),
    ]
    cols = 2
    card_w = Inches(6.0); card_h = Inches(2.0); gap = Inches(0.10)
    start_x = Inches(0.6); start_y = Inches(2.2)
    for i, (head, body, accent) in enumerate(metrics):
        r = i // cols; c = i % cols
        x = start_x + (card_w + gap) * c
        y = start_y + (card_h + gap) * r
        add_rounded_rect(s, x, y, card_w, card_h, fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
        add_rect(s, x, y, Inches(0.10), card_h, fill=accent)
        add_textbox(s, x + Inches(0.35), y + Inches(0.30), card_w - Inches(0.5), Inches(0.45),
                    head, size=18, bold=True, color=SLATE_900)
        add_textbox(s, x + Inches(0.35), y + Inches(0.85), card_w - Inches(0.5), Inches(0.9),
                    body, size=14, color=SLATE_500, line_spacing=1.4)

    add_rounded_rect(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.55),
                     fill=NAVY_SOFT, line=None)
    add_textbox(s, Inches(0.85), Inches(6.55), Inches(11.7), Inches(0.4),
        'Tooling: Ragas · TruLens · DeepEval · or 30 lines of pytest. Just measure something.',
        size=13, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "Evaluation. The thing nobody wants to do, the thing that separates "
        "real RAG systems from demos. Without metrics, you can't tell if "
        "yesterday's prompt change made things better or worse. Track four "
        "categories. One — retrieval quality. Precision-at-K asks 'of the K "
        "passages we retrieved, how many were relevant?' Recall-at-K asks "
        "'of all relevant passages that exist, how many did we find?' MRR — "
        "Mean Reciprocal Rank — is great when there's typically one right "
        "answer. Two — faithfulness. Does the generated answer actually "
        "match the context, or did the model wander off and hallucinate? "
        "This is usually scored by another LLM. Three — answer relevance. "
        "Does the answer actually address the question, even if it's "
        "faithful? Four — latency and cost. P95 milliseconds, tokens per "
        "query, dollars per hundred queries. Tools — Ragas, TruLens, "
        "DeepEval. Or you can write thirty lines of pytest with a hand-curated "
        "test set. The point isn't tooling sophistication — the point is "
        "you measure SOMETHING so changes are comparable. If you take one "
        "thing from this video, it's: build the eval harness early.")
    return s


def s17_production(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Operations',
                   title='Production patterns')

    cards = [
        ('Cache aggressively',  'Hash query+filters → store answer · most user queries repeat', CYAN),
        ('Citations always',    'Every answer carries the source IDs used · UI links back',     EMERALD),
        ('Observability',       'Log: query, retrieved IDs, prompt, answer, latency, cost',     AMBER),
        ('Guardrails',          'PII redaction · max context · prompt-injection filtering',     ROSE),
    ]
    cols = 2
    card_w = Inches(6.0); card_h = Inches(2.0); gap = Inches(0.10)
    start_x = Inches(0.6); start_y = Inches(2.2)
    for i, (head, body, accent) in enumerate(cards):
        r = i // cols; c = i % cols
        x = start_x + (card_w + gap) * c
        y = start_y + (card_h + gap) * r
        add_rounded_rect(s, x, y, card_w, card_h, fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
        add_rect(s, x, y, Inches(0.10), card_h, fill=accent)
        add_textbox(s, x + Inches(0.35), y + Inches(0.25), card_w - Inches(0.5), Inches(0.4),
                    head, size=18, bold=True, color=SLATE_900)
        add_textbox(s, x + Inches(0.35), y + Inches(0.75), card_w - Inches(0.5), Inches(1.2),
                    body, size=13, color=SLATE_700, line_spacing=1.4)
    add_speaker_notes(s,
        "Production patterns. Four things you absolutely need before going "
        "live. One — cache aggressively. The vast majority of user queries "
        "repeat with minor variation. Hash the normalized query plus filters, "
        "store the answer, return it on hit. You'll cut LLM bills by 70 "
        "percent or more. Two — citations always, never optional. Every "
        "answer carries the source IDs it used; the UI links each citation "
        "back to the underlying record. This single pattern transforms user "
        "trust — they go from 'is this right?' to 'let me click and verify.' "
        "Three — observability. Log the query, the retrieved IDs, the full "
        "prompt, the answer, the latency, the cost — for every single "
        "request. Without this you cannot debug RAG in production; the "
        "system is too non-deterministic. Four — guardrails. PII redaction "
        "before logging, max context token caps, prompt-injection filtering "
        "on user input. The injection threat is real — users will try to "
        "trick the model into ignoring its instructions. Skip these patterns "
        "and you'll have a great demo and a bad production incident.")
    return s


def s18_tools(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Reference',
                   title='Tools we’ll use across the series')

    cols = [
        ('Models',     ['OpenAI · Anthropic', 'Llama 3 · Mistral'],   INDIGO),
        ('Embeddings', ['text-embedding-3', 'BGE / E5'],              CYAN),
        ('Stores',     ['pgvector', 'Pinecone · Weaviate'],           EMERALD),
        ('Orchestrate',['LangChain', 'LlamaIndex · plain Python'],    AMBER),
        ('Evaluate',   ['Ragas', 'TruLens · DeepEval'],               ROSE),
    ]
    n = len(cols)
    card_w = Inches(2.40); gap = Inches(0.12)
    total_w = card_w*n + gap*(n-1)
    start_x = (SLIDE_W - total_w) // 2
    y = Inches(2.6); h = Inches(3.6)
    for i, (head, items, accent) in enumerate(cols):
        x = start_x + (card_w + gap) * i
        add_rounded_rect(s, x, y, card_w, h, fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
        add_rect(s, x, y, card_w, Inches(0.55), fill=accent)
        add_textbox(s, x, y + Inches(0.10), card_w, Inches(0.4),
                    head.upper(), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            add_textbox(s, x + Inches(0.2), y + Inches(0.85) + Inches(0.55)*j,
                        card_w - Inches(0.4), Inches(0.5),
                        f'• {item}', size=12, color=SLATE_700)
    add_textbox(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
                'Pick one per column. Ship. Iterate.',
                size=14, color=SLATE_500, align=PP_ALIGN.CENTER, bold=True)
    add_speaker_notes(s,
        "Tools cheat sheet for the rest of the series. Models — pick OpenAI "
        "or Anthropic for the easiest starting point; switch to Llama 3 or "
        "Mistral when you need open-weights or on-prem. Embeddings — OpenAI's "
        "text-embedding-3 is fine for general purpose; BGE or E5 if you have "
        "a domain corpus. Stores — pgvector first, almost always. Move to "
        "Pinecone or Weaviate when you actually have to. Orchestration — "
        "LangChain has the biggest ecosystem and the most footguns; "
        "LlamaIndex is more focused on retrieval; plain Python is honestly "
        "fine for most projects and gives you the cleanest debugging story. "
        "Evaluation — Ragas to start, TruLens or DeepEval when you need "
        "more metrics. The meta-rule: pick one per column and SHIP. Do not "
        "spend a week comparing every option before you build your first "
        "version. Start simple, ship, measure, then add complexity only when "
        "your metrics force you to. Speed of iteration beats sophistication "
        "of stack every time.")
    return s


def s19_roadmap(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Coming up',
                   title='Episodes ahead')

    eps = [
        ('Ep 02', 'Embeddings deep dive — choosing the right model',        INDIGO),
        ('Ep 03', 'Vector stores benchmarked — pgvector vs managed',         CYAN),
        ('Ep 04', 'Banking — sanctions screening RAG (OFAC + AML)',          EMERALD),
        ('Ep 05', 'Healthcare — querying clinical notes safely',             AMBER),
        ('Ep 06', 'Evaluation — building the harness',                       ROSE),
    ]
    y = Inches(2.4); h = Inches(0.75); gap = Inches(0.12)
    for i, (label, body, accent) in enumerate(eps):
        ry = y + (h + gap) * i
        add_rounded_rect(s, Inches(0.6), ry, Inches(12.1), h, fill=WHITE, line=SLATE_200, line_w=Pt(0.5))
        add_rounded_rect(s, Inches(0.85), ry + Inches(0.18), Inches(1.6), Inches(0.42),
                         fill=accent, line=None)
        add_textbox(s, Inches(0.85), ry + Inches(0.21), Inches(1.6), Inches(0.4),
                    label.upper(), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(2.65), ry + Inches(0.20), Inches(9.8), Inches(0.45),
                    body, size=15, color=SLATE_900)
    add_speaker_notes(s,
        "Quick look at what's coming next. Episode two — embeddings deep "
        "dive. We'll compare OpenAI ada-3, the various BGE flavours, and "
        "open-source sentence-transformers, and I'll show you the actual "
        "quality difference on a real corpus. Episode three — vector stores "
        "benchmarked side by side. pgvector versus Pinecone versus Weaviate "
        "on the same data, the same queries, with timings and costs. Then "
        "we shift into vertical case studies. Episode four — banking "
        "sanctions screening with RAG, working with public OFAC data. "
        "Episode five — healthcare, querying clinical notes safely with "
        "PII redaction and access control. Episode six — the one most "
        "channels skip — building the evaluation harness, because if you "
        "remember nothing else from this series, remember that you can't "
        "improve what you can't measure. Subscribe so you don't miss any of "
        "these, and tell me in the comments which one you want me to "
        "prioritize.")
    return s


def s20_outro(prs, _total):
    s = base_slide(prs, no_chrome=True)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY_DARK)
    add_rect(s, 0, 0, Inches(7.0), SLIDE_H, fill=NAVY)
    add_rect(s, 0, Inches(7.30), SLIDE_W, Inches(0.20), fill=CYAN)
    add_rect(s, 0, Inches(7.20), Inches(5.0), Inches(0.10), fill=AMBER)

    add_textbox(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
                'THANKS FOR WATCHING', size=14, bold=True, color=CYAN)
    add_textbox(s, Inches(0.8), Inches(2.3), Inches(11), Inches(1.4),
                'Like · Subscribe · Comment', size=50, bold=True, color=WHITE)
    add_textbox(s, Inches(0.8), Inches(3.85), Inches(11), Inches(0.6),
                'Tell me which vertical I should ship next:',
                size=20, color=CYAN_SOFT)
    add_textbox(s, Inches(0.8), Inches(4.45), Inches(11), Inches(0.6),
                'Banking · Healthcare · BioMed · Auto · Education · Search',
                size=18, color=SLATE_200)

    add_rounded_rect(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.0),
                     fill=NAVY_DARK, line=CYAN, line_w=Pt(1.0))
    add_textbox(s, Inches(1.0), Inches(5.65), Inches(11.0), Inches(0.4),
                'WORKED EXAMPLE — SOURCE CODE', size=10, bold=True, color=CYAN)
    add_textbox(s, Inches(1.0), Inches(5.95), Inches(11.0), Inches(0.5),
                'github.com/javakishore-veleti/FX-Trade-Analytics-AWs-OpenSearch',
                size=14, color=WHITE, font='Menlo')

    add_textbox(s, Inches(0.8), Inches(6.85), Inches(11), Inches(0.4),
                'Channel: Training My Hobby · Series: RAG Mastery',
                size=11, color=SLATE_400)
    add_speaker_notes(s,
        "And that's RAG Mastery, episode one. If this gave you a clearer "
        "mental model of RAG, smash the like button — it genuinely helps "
        "the channel. Subscribe so you don't miss the embeddings deep dive "
        "next week. Drop a comment with which vertical you want me to ship "
        "first — banking, healthcare, biomed, auto, education, or search. "
        "And the full source code for the FX trading worked example I "
        "showed today is in the description and on the slide here on screen "
        "— go fork it, break it, build something on top of it. Thanks for "
        "watching. See you in episode two.")
    return s


# ---------------------------------------------------------------- main
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        (s01_title,         False),
        (s02_rag_overview,  True),
        (s03_rag_buckets,   True),
        (s04_project_arch,  True),
        (s05_bridge,        True),
        (s06_what_is_rag,   True),
        (s07_pipeline,      True),
        (s08_embeddings,    True),
        (s09_vector_store,  True),
        (s10_retrieval,     True),
        (s11_augment,       True),
        (s12_generate,      True),
        (s13_rag_over_fx,   True),
        (s14_hybrid,        True),
        (s15_chunking,      True),
        (s16_eval,          True),
        (s17_production,    True),
        (s18_tools,         True),
        (s19_roadmap,       True),
        (s20_outro,         False),
    ]
    total = len(builders)
    for i, (build, takes_page) in enumerate(builders, 1):
        if takes_page:
            build(prs, i, total)
        else:
            build(prs, total)

    out = Path(__file__).parent / '01-RAG-and-FX-Trade-Analytics.pptx'
    prs.save(out)
    print(f'Wrote {out}')
    print(f'Slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()
