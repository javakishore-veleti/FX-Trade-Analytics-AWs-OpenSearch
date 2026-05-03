"""Generate a 15-18 minute YouTube deck centered on the AWS announcement
"OpenSearch UI supports cross-region data access to OpenSearch domains"
(May 1, 2026), using the FX Trade Analytics project in this repo as the
worked example.

Run:
    python3 build_aws_opensearch_xregion_deck.py

Output:
    AWS-OpenSearch-UI-Cross-Region-FX-Demo.pptx (alongside this script)

Source link cited on slide 2:
    https://aws.amazon.com/about-aws/whats-new/2026/05/opensearch-ui-cross-region-data-access-domains/
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pathlib import Path

# ---------------------------------------------------------------- design system
NAVY         = RGBColor(0x1E, 0x3A, 0x8A)
NAVY_DARK    = RGBColor(0x0F, 0x2A, 0x5C)
NAVY_SOFT    = RGBColor(0xDB, 0xEA, 0xFE)
INDIGO       = RGBColor(0x31, 0x2E, 0x81)
INDIGO_SOFT  = RGBColor(0xE0, 0xE7, 0xFF)
CYAN         = RGBColor(0x06, 0xB6, 0xD4)
CYAN_SOFT    = RGBColor(0xCF, 0xFA, 0xFE)
EMERALD      = RGBColor(0x10, 0xB9, 0x81)
EMERALD_SOFT = RGBColor(0xD1, 0xFA, 0xE5)
AMBER        = RGBColor(0xF5, 0x9E, 0x0B)
ROSE         = RGBColor(0xF4, 0x3F, 0x5E)
AWS_ORANGE   = RGBColor(0xFF, 0x99, 0x00)   # AWS brand orange
AWS_DARK     = RGBColor(0x23, 0x2F, 0x3E)   # AWS brand dark

WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
SLATE_50     = RGBColor(0xF8, 0xFA, 0xFC)
SLATE_100    = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_200    = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_400    = RGBColor(0x94, 0xA3, 0xB8)
SLATE_500    = RGBColor(0x64, 0x74, 0x8B)
SLATE_700    = RGBColor(0x33, 0x41, 0x55)
SLATE_900    = RGBColor(0x0F, 0x17, 0x2A)

FONT = 'Inter'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

CHANNEL = 'Training My Hobby · Kishore Veleti'
SERIES  = "AWS What's New · OpenSearch UI Cross-Region"

# ---------------------------------------------------------------- helpers
def set_text(tf, text, *, size=18, bold=False, color=SLATE_700, align=PP_ALIGN.LEFT,
             font=FONT, line_spacing=1.15):
    tf.word_wrap = True
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0);  tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    if not p.runs:
        r = p.add_run()
    else:
        r = p.runs[0]
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, **kwargs):
    box = slide.shapes.add_textbox(x, y, w, h)
    set_text(box.text_frame, text, **kwargs)
    return box


def add_bullets(slide, x, y, w, h, items, *, size=18, color=SLATE_700,
                bullet_color=CYAN, line_spacing=1.4, gap_after=Pt(8)):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0);  tf.margin_bottom = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = gap_after
        rb = p.add_run()
        rb.text = '●  '
        rb.font.name = FONT
        rb.font.size = Pt(size)
        rb.font.color.rgb = bullet_color
        rt = p.add_run()
        rt.text = item
        rt.font.name = FONT
        rt.font.size = Pt(size)
        rt.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, *, fill=WHITE, line=None, line_w=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_rounded_rect(slide, x, y, w, h, *, fill=WHITE, line=None, line_w=Pt(0), radius=0.10):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = radius
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_arrow(slide, x, y, w, h, *, fill=CYAN):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False
    return s


def base_slide(prs, *, page=None, total=None, title=None, eyebrow=None, no_chrome=False):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    if no_chrome:
        return slide
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.10), fill=NAVY)
    add_rect(slide, 0, Inches(0.10), Inches(4.5), Inches(0.04), fill=AWS_ORANGE)
    add_rect(slide, Inches(4.5), Inches(0.10), Inches(2.0), Inches(0.04), fill=CYAN)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
                CHANNEL, size=10, color=SLATE_400)
    add_textbox(slide, Inches(0.5), Inches(7.20), Inches(8), Inches(0.3),
                SERIES, size=10, color=SLATE_500, bold=True)
    if page is not None:
        add_textbox(slide, Inches(12.0), Inches(7.10), Inches(1.0), Inches(0.3),
                    f'{page:02d} / {total:02d}', size=10, color=SLATE_400, align=PP_ALIGN.RIGHT)
    if eyebrow:
        add_textbox(slide, Inches(0.6), Inches(0.45), Inches(10), Inches(0.4),
                    eyebrow.upper(), size=11, bold=True, color=AWS_ORANGE)
    if title:
        add_textbox(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.95),
                    title, size=32, bold=True, color=SLATE_900)
        add_rect(slide, Inches(0.6), Inches(1.78), Inches(0.5), Inches(0.06), fill=AWS_ORANGE)
    return slide


def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ============================================================== slides
def s01_title(prs, _total):
    s = base_slide(prs, no_chrome=True)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=AWS_DARK)
    add_rect(s, 0, 0, Inches(7.5), SLIDE_H, fill=NAVY_DARK)
    add_rect(s, 0, Inches(7.30), SLIDE_W, Inches(0.20), fill=AWS_ORANGE)
    add_rect(s, 0, Inches(7.20), Inches(5.0), Inches(0.10), fill=CYAN)

    add_textbox(s, Inches(0.8), Inches(2.0), Inches(11), Inches(0.4),
                "AWS WHAT'S NEW · MAY 2026", size=14, bold=True, color=AWS_ORANGE)
    add_textbox(s, Inches(0.8), Inches(2.5), Inches(12), Inches(1.5),
                'OpenSearch UI', size=68, bold=True, color=WHITE)
    add_textbox(s, Inches(0.8), Inches(3.85), Inches(12), Inches(0.95),
                'Cross-region data access — without moving data',
                size=28, color=CYAN_SOFT)
    add_textbox(s, Inches(0.8), Inches(4.85), Inches(12), Inches(0.6),
                'Worked example: a real-time multi-region FX trading platform',
                size=18, color=SLATE_200)
    add_textbox(s, Inches(0.8), Inches(5.5), Inches(12), Inches(0.5),
                '~15 minute walkthrough · architecture, demo, why it matters',
                size=14, color=SLATE_400)
    add_textbox(s, Inches(0.8), Inches(6.5), Inches(12), Inches(0.4),
                'Hosted by Kishore Veleti', size=12, color=SLATE_400, bold=True)

    add_speaker_notes(s,
        "Welcome back. On May 1st, 2026, AWS quietly shipped a feature that "
        "matters a lot more than its single-line announcement suggests — "
        "cross-region data access for OpenSearch UI. In the next fifteen "
        "minutes, I'll walk you through what it is, why it's a big deal for "
        "anyone running OpenSearch in more than one region, and I'll "
        "demonstrate it on a real project I built — a multi-region FX trade "
        "analytics platform whose architecture happens to be perfectly "
        "positioned for this. By the end you'll know if this feature is "
        "useful for you and exactly how to try it. Let's go.")
    return s


def s02_announcement(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow="AWS what's new · May 1, 2026",
                   title='The announcement, in two sentences')

    # AWS-style card
    add_rounded_rect(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(2.4),
                     fill=AWS_DARK, line=None)
    add_textbox(s, Inches(0.95), Inches(2.4), Inches(11.5), Inches(0.4),
                'OFFICIAL AWS RELEASE', size=11, bold=True, color=AWS_ORANGE)
    add_textbox(s, Inches(0.95), Inches(2.75), Inches(11.5), Inches(0.7),
                'OpenSearch UI supports cross-region data access to OpenSearch domains',
                size=22, bold=True, color=WHITE, line_spacing=1.3)
    add_textbox(s, Inches(0.95), Inches(3.6), Inches(11.5), Inches(0.95),
        '“Access OpenSearch domains hosted in different AWS Regions from within '
        'a single OpenSearch UI application — without switching endpoints or '
        'replicating data.”',
        size=14, color=CYAN_SOFT, line_spacing=1.5)

    # Quick facts row
    facts = [
        ('Date',   'May 1, 2026'),
        ('Scope',  'Public + VPC domains'),
        ('Auth',   'IAM · IAM Identity Center'),
        ('Combo',  'Plus cross-account access'),
    ]
    for i, (k, v) in enumerate(facts):
        x = Inches(0.6) + Inches(3.05) * i
        add_rounded_rect(s, x, Inches(4.95), Inches(2.95), Inches(1.0),
                         fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
        add_textbox(s, x + Inches(0.20), Inches(5.05), Inches(2.55), Inches(0.3),
                    k.upper(), size=10, bold=True, color=AWS_ORANGE)
        add_textbox(s, x + Inches(0.20), Inches(5.35), Inches(2.55), Inches(0.55),
                    v, size=13, color=SLATE_900, line_spacing=1.3)

    add_textbox(s, Inches(0.6), Inches(6.25), Inches(12.1), Inches(0.4),
        'aws.amazon.com/about-aws/whats-new/2026/05/opensearch-ui-cross-region-data-access-domains/',
        size=11, color=SLATE_500, font='Menlo')

    add_speaker_notes(s,
        "Here's the announcement, verbatim. May 1st 2026. OpenSearch UI now "
        "supports cross-region data access — meaning you can query OpenSearch "
        "domains in multiple regions from a single UI application, without "
        "switching endpoints, without replicating data. Four quick facts: "
        "available immediately, works on both public and VPC domains, "
        "supports IAM and IAM Identity Center for end-user auth, and combines "
        "with the cross-account feature that launched earlier this year. The "
        "URL on screen is the official AWS what's-new page — go read it after "
        "the video.")
    return s


def s03_problem(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Why this matters',
                   title='The trade-off this collapses')

    # Two pillars: residency vs visibility
    add_rounded_rect(s, Inches(0.6), Inches(2.2), Inches(5.95), Inches(4.3),
                     fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
    add_rect(s, Inches(0.6), Inches(2.2), Inches(0.10), Inches(4.3), fill=ROSE)
    add_textbox(s, Inches(0.95), Inches(2.4), Inches(5.5), Inches(0.4),
                'PRESSURE 1 · DATA RESIDENCY', size=11, bold=True, color=ROSE)
    add_textbox(s, Inches(0.95), Inches(2.75), Inches(5.5), Inches(0.5),
                'Keep data in the region where it was born',
                size=18, bold=True, color=SLATE_900)
    add_bullets(s, Inches(0.95), Inches(3.45), Inches(5.4), Inches(3.0),
                ['GDPR · MAS · RBI · regional compliance regimes',
                 'Customer contractual requirements',
                 'Latency for regional users',
                 'Inter-region egress is expensive'],
                size=13, color=SLATE_700, gap_after=Pt(8))

    add_rounded_rect(s, Inches(6.85), Inches(2.2), Inches(5.95), Inches(4.3),
                     fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
    add_rect(s, Inches(6.85), Inches(2.2), Inches(0.10), Inches(4.3), fill=EMERALD)
    add_textbox(s, Inches(7.20), Inches(2.4), Inches(5.5), Inches(0.4),
                'PRESSURE 2 · CENTRAL VISIBILITY', size=11, bold=True, color=EMERALD)
    add_textbox(s, Inches(7.20), Inches(2.75), Inches(5.5), Inches(0.5),
                'Analysts need one pane of glass',
                size=18, bold=True, color=SLATE_900)
    add_bullets(s, Inches(7.20), Inches(3.45), Inches(5.4), Inches(3.0),
                ['Global trading book oversight',
                 'Cross-region anomaly detection',
                 'Centralized observability dashboards',
                 'One UI, one auth, one query language'],
                size=13, color=SLATE_700, gap_after=Pt(8), bullet_color=EMERALD)

    add_rounded_rect(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4),
                     fill=AWS_DARK, line=None)
    add_textbox(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.4),
        'For years you had to pick one. Now you don’t.',
        size=14, color=AWS_ORANGE, bold=True, align=PP_ALIGN.CENTER)

    add_speaker_notes(s,
        "Why does cross-region matter? Because every team running OpenSearch "
        "globally has been stuck between two pressures. On the left — data "
        "residency. GDPR, the Monetary Authority of Singapore, India's RBI, "
        "customer contractual rules — all push you to keep data in its origin "
        "region. Plus the practical pressures: latency for local users, and "
        "inter-region egress fees that add up fast. On the right — central "
        "visibility. Analysts need one pane of glass to see global trading "
        "books, run anomaly detection across regions, observe the platform "
        "as a whole. For years you had to pick one — replicate data and "
        "blow your egress bill, or accept fragmented dashboards. This feature "
        "collapses the trade-off.")
    return s


def s04_before(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Before this feature',
                   title="What teams used to do (and pay for)")

    options = [
        ('Cross-cluster replication',
         'Replicate every region into a central one. Doubles storage, blows up egress, breaks residency.',
         ROSE),
        ('Endpoint juggling',
         'Five bookmarks for five regions. Switch contexts, copy-paste queries, no joins.',
         AMBER),
        ('Custom federated proxy',
         'Build a query router yourself. New surface area to maintain and authenticate.',
         INDIGO),
        ('Just give up',
         'Per-region dashboards · no global view · the most common outcome',
         SLATE_500),
    ]
    cols = 2
    card_w = Inches(6.0); card_h = Inches(2.0); gap = Inches(0.10)
    start_x = Inches(0.6); start_y = Inches(2.2)
    for i, (head, body, accent) in enumerate(options):
        r = i // cols; c = i % cols
        x = start_x + (card_w + gap) * c
        y = start_y + (card_h + gap) * r
        add_rounded_rect(s, x, y, card_w, card_h, fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
        add_rect(s, x, y, Inches(0.10), card_h, fill=accent)
        add_textbox(s, x + Inches(0.35), y + Inches(0.25), card_w - Inches(0.5), Inches(0.4),
                    head, size=17, bold=True, color=SLATE_900)
        add_textbox(s, x + Inches(0.35), y + Inches(0.75), card_w - Inches(0.5), Inches(1.2),
                    body, size=13, color=SLATE_700, line_spacing=1.4)

    add_speaker_notes(s,
        "Before May 1st, you had four options, and all four are bad. Option "
        "one — cross-cluster replication. Pull every region's data into a "
        "central region for analysis. You double storage costs, blow up "
        "egress, and the moment you replicate data across borders, your "
        "compliance team starts asking hard questions. Option two — endpoint "
        "juggling. Five regional dashboards, five bookmarks, copy-paste "
        "queries, no way to join. Option three — build your own federated "
        "query proxy. Now you own a new piece of infrastructure and a new "
        "auth surface area. Option four, the most common in the wild — just "
        "give up. Per-region dashboards, no global view. Pick your poison.")
    return s


def s05_after(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='After this feature',
                   title='One UI · many regions · no movement')

    # Visual: 3 region domains feeding into a single UI
    # left side - 3 region cards
    regions = [
        ('us-east-1',   'Domain', NAVY),
        ('eu-west-1',   'Domain', CYAN),
        ('ap-south-1',  'Domain', EMERALD),
    ]
    for i, (r, sub, accent) in enumerate(regions):
        y = Inches(2.4) + Inches(1.4) * i
        add_rounded_rect(s, Inches(0.7), y, Inches(2.8), Inches(1.1),
                         fill=WHITE, line=accent, line_w=Pt(1.5))
        add_textbox(s, Inches(0.85), y + Inches(0.18), Inches(2.5), Inches(0.4),
                    'OPENSEARCH', size=10, bold=True, color=accent)
        add_textbox(s, Inches(0.85), y + Inches(0.45), Inches(2.5), Inches(0.5),
                    r, size=18, bold=True, color=SLATE_900, font='Menlo')
        add_textbox(s, Inches(0.85), y + Inches(0.78), Inches(2.5), Inches(0.3),
                    sub, size=11, color=SLATE_500)
        # arrow pointing right
        add_arrow(s, Inches(3.7), y + Inches(0.4), Inches(0.6), Inches(0.30), fill=accent)

    # Center / right: single UI box
    add_rounded_rect(s, Inches(4.6), Inches(2.6), Inches(8.1), Inches(4.0),
                     fill=NAVY_DARK, line=AWS_ORANGE, line_w=Pt(2.0))
    add_textbox(s, Inches(4.85), Inches(2.85), Inches(7.5), Inches(0.4),
                'OPENSEARCH UI · SINGLE APPLICATION', size=11, bold=True, color=AWS_ORANGE)
    add_textbox(s, Inches(4.85), Inches(3.25), Inches(7.5), Inches(0.6),
                'One dashboard.', size=24, bold=True, color=WHITE)
    add_textbox(s, Inches(4.85), Inches(3.85), Inches(7.5), Inches(0.6),
                'One auth. One query language.', size=18, color=CYAN_SOFT)
    add_textbox(s, Inches(4.85), Inches(4.6), Inches(7.5), Inches(1.7),
        '✓  Data stays in its origin region — zero replication\n'
        '✓  Zero inter-region egress for storage\n'
        '✓  Per-region latency + availability preserved\n'
        '✓  Compliant with regional data residency rules',
        size=13, color=SLATE_200, line_spacing=1.7)

    add_speaker_notes(s,
        "Here's what the world looks like after May 1st. Three regional "
        "OpenSearch domains on the left — each holding its local trades, "
        "subject to its local rules. On the right, a single OpenSearch UI "
        "application that talks to all three. The user sees one dashboard, "
        "logs in once, runs one query language. But — and this is the part "
        "that matters — no data moves. The fx-trades-us-east-1 index stays "
        "in us-east-1. The fx-trades-eu-west-1 index stays in eu-west-1. The "
        "UI federates the queries at runtime. You get global visibility "
        "without giving up any of the four properties on the right.")
    return s


def s06_architecture(prs, page, total):
    """Single slide that introduces the project AND shows the multi-region
    deployment topology — OpenSearch UI federating across regions."""
    s = base_slide(prs, page=page, total=total,
                   eyebrow='The worked example',
                   title='Meet FX Trade Analytics — a reference implementation')

    # Subtitle / framing
    add_textbox(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.5),
        'Open-source reference implementation showcasing the new AWS OpenSearch '
        'cross-region UI feature on a real, production-shaped workload.',
        size=14, color=SLATE_700, line_spacing=1.4)

    # ── OpenSearch UI hero box (top, central, AWS-branded) ──
    ui_w = Inches(7.0); ui_h = Inches(0.95)
    ui_x = (SLIDE_W - ui_w) // 2
    ui_y = Inches(2.65)
    add_rounded_rect(s, ui_x, ui_y, ui_w, ui_h, fill=AWS_DARK,
                     line=AWS_ORANGE, line_w=Pt(2.0))
    add_textbox(s, ui_x, ui_y + Inches(0.12), ui_w, Inches(0.32),
                'AWS · OPENSEARCH UI', size=11, bold=True, color=AWS_ORANGE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, ui_x, ui_y + Inches(0.42), ui_w, Inches(0.4),
                'Single application — federates queries across regions + accounts',
                size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Region columns (3 deployments, parallel) ──
    regions = [
        ('us-east-1',  'N. Virginia', NAVY),
        ('eu-west-1',  'Ireland',     CYAN),
        ('ap-south-1', 'Mumbai',      EMERALD),
    ]
    col_w = Inches(3.93)
    gap   = Inches(0.165)
    col_y = Inches(4.10)
    col_h = Inches(2.78)

    # Connector lines from UI box bottom down to each region top
    ui_bottom_y = ui_y + ui_h
    for i, (_, _, accent) in enumerate(regions):
        col_x = Inches(0.6) + (col_w + gap) * i
        cx = col_x + col_w / 2
        # vertical connector
        add_rect(s, cx - Inches(0.012), ui_bottom_y, Inches(0.024),
                 col_y - ui_bottom_y, fill=AWS_ORANGE)
        # arrowhead at the bottom
        add_arrow(s, cx - Inches(0.10), col_y - Inches(0.18),
                  Inches(0.20), Inches(0.18), fill=AWS_ORANGE)

    # Each region: top-down stack showing the application tiers
    for i, (region, city, accent) in enumerate(regions):
        col_x = Inches(0.6) + (col_w + gap) * i
        # Outer container
        add_rounded_rect(s, col_x, col_y, col_w, col_h, fill=WHITE,
                         line=accent, line_w=Pt(1.5))
        # Region header strip
        add_rect(s, col_x, col_y, col_w, Inches(0.42), fill=accent)
        add_textbox(s, col_x, col_y + Inches(0.05), col_w, Inches(0.22),
                    region, size=14, bold=True, color=WHITE,
                    font='Menlo', align=PP_ALIGN.CENTER)
        add_textbox(s, col_x, col_y + Inches(0.24), col_w, Inches(0.18),
                    city, size=10, color=WHITE, align=PP_ALIGN.CENTER)

        # Inner tiers, top-down
        inner_x = col_x + Inches(0.18)
        inner_w = col_w - Inches(0.36)
        ty = col_y + Inches(0.55)
        tier_h = Inches(0.50)
        tier_gap = Inches(0.05)

        tiers = [
            ('WEB TIER',           'Customer Portal · Admin Portal', NAVY_SOFT, NAVY_DARK),
            ('APPLICATION TIER',   'Trade · Risk · Indexer · Master Data', INDIGO_SOFT, NAVY_DARK),
            ('STREAMING + STATE',  'Kafka · Postgres', AMBER, WHITE),
        ]
        for label, body, fill, fg in tiers:
            add_rounded_rect(s, inner_x, ty, inner_w, tier_h, fill=fill, line=None)
            add_textbox(s, inner_x, ty + Inches(0.04), inner_w, Inches(0.20),
                        label, size=8, bold=True, color=fg, align=PP_ALIGN.CENTER)
            add_textbox(s, inner_x, ty + Inches(0.23), inner_w, Inches(0.27),
                        body, size=10, color=fg, align=PP_ALIGN.CENTER)
            ty += tier_h + tier_gap

        # OpenSearch domain tier — highlighted, this is what UI federates
        os_h = Inches(0.55)
        add_rounded_rect(s, inner_x, ty, inner_w, os_h,
                         fill=AWS_DARK, line=AWS_ORANGE, line_w=Pt(1.0))
        add_textbox(s, inner_x, ty + Inches(0.05), inner_w, Inches(0.20),
                    'OPENSEARCH DOMAIN', size=8, bold=True, color=AWS_ORANGE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, inner_x, ty + Inches(0.25), inner_w, Inches(0.28),
                    f'fx-trades-{region}', size=10, color=WHITE,
                    font='Menlo', align=PP_ALIGN.CENTER)

    add_speaker_notes(s,
        "This single slide is the introduction and the architecture. FX "
        "Trade Analytics is an open-source reference implementation — "
        "deliberately built to showcase exactly this new AWS feature on a "
        "real, production-shaped workload, not a toy. Look at the picture: "
        "every region — us-east-1, eu-west-1, ap-south-1 — runs an identical "
        "deployment of the entire stack. The web tier with our two Angular "
        "portals. The application tier with our four Spring Boot services — "
        "trade, risk, indexer, and master data. Streaming and state — Kafka "
        "for events, Postgres for reference data. And at the bottom of each "
        "region, the OpenSearch domain holding that region's trade index. "
        "What's new is the orange piece at the top — OpenSearch UI as a "
        "single application, federating queries across all three regional "
        "domains. Same architecture, three regions, one dashboard. From "
        "here on, every concept lands on this picture.")
    return s


def s08_partitioned(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='The design choice that pays off',
                   title='Region-partitioned indices')

    add_rounded_rect(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.4),
                     fill=NAVY_SOFT, line=None)
    add_textbox(s, Inches(0.95), Inches(2.4), Inches(11.5), Inches(0.4),
                'INDEX NAMING CONVENTION', size=11, bold=True, color=NAVY_DARK)
    add_textbox(s, Inches(0.95), Inches(2.75), Inches(11.5), Inches(0.7),
                'fx-trades-{region}',
                size=28, bold=True, color=NAVY_DARK, font='Menlo')

    # 3 examples
    examples = [
        ('fx-trades-us-east-1', 'NYC trades', NAVY),
        ('fx-trades-eu-west-1', 'EU trades',  CYAN),
        ('fx-trades-ap-south-1','APAC trades', EMERALD),
    ]
    y = Inches(3.95)
    for i, (idx, sub, accent) in enumerate(examples):
        x = Inches(0.6) + Inches(4.05) * i
        add_rounded_rect(s, x, y, Inches(3.95), Inches(1.4),
                         fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
        add_rect(s, x, y, Inches(3.95), Inches(0.10), fill=accent)
        add_textbox(s, x + Inches(0.2), y + Inches(0.30), Inches(3.6), Inches(0.5),
                    idx, size=15, bold=True, color=SLATE_900, font='Menlo')
        add_textbox(s, x + Inches(0.2), y + Inches(0.85), Inches(3.6), Inches(0.4),
                    sub, size=13, color=SLATE_500)

    # Bottom callout
    add_rounded_rect(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.0),
                     fill=AWS_DARK, line=AWS_ORANGE, line_w=Pt(1.0))
    add_textbox(s, Inches(0.95), Inches(5.85), Inches(11.5), Inches(0.4),
                'WHY THIS MATTERS NOW', size=11, bold=True, color=AWS_ORANGE)
    add_textbox(s, Inches(0.95), Inches(6.15), Inches(11.5), Inches(0.85),
        'OpenSearch UI’s federated query routes by index pattern. fx-trades-* '
        'across regions just works.',
        size=14, color=WHITE, line_spacing=1.4)

    add_speaker_notes(s,
        "The architectural decision that makes this whole thing click: every "
        "trade is indexed with a region-suffixed index name — fx-trades-us-east-1, "
        "fx-trades-eu-west-1, fx-trades-ap-south-1. One index per region. "
        "Each one lives in the OpenSearch domain in its own region. Why "
        "this matters now: OpenSearch UI's cross-region federated query "
        "routes by index pattern. So a query against fx-trades-* "
        "automatically hits all three regional domains. No code changes. "
        "No data movement. Just a wildcard.")
    return s


def s09_forward_looking(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Lessons',
                   title='We did this before AWS shipped the feature')

    add_textbox(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(0.6),
        'Region-partitioning isn’t about anticipating AWS roadmaps — it’s a structural choice.',
        size=16, color=SLATE_700, line_spacing=1.4)

    items = [
        'Data residency is a forcing function — it’s rarely optional anyway',
        'Per-region scaling, sharding, and rollouts are easier when storage is partitioned',
        'Region-as-tenant gives you natural blast-radius isolation',
        'When AWS ships a federation feature, your data is already in the right shape',
    ]
    add_bullets(s, Inches(0.6), Inches(3.1), Inches(12), Inches(3.5),
                items, size=18, color=SLATE_900, gap_after=Pt(14))

    add_textbox(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.4),
        'Design for residency. Centralization is a feature you can buy later.',
        size=14, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "There's a lesson here. We didn't partition by region because we "
        "knew this AWS feature was coming — we partitioned because residency "
        "is a forcing function for almost every regulated industry. But "
        "what's interesting is that designing for residency tends to leave "
        "you better positioned for everything else: per-region scaling, "
        "blast-radius isolation, regional rollouts. And when AWS finally "
        "ships a federation feature, your data is already in the shape it "
        "needs to be in. Generalize the rule: design for residency, treat "
        "centralization as something you buy later.")
    return s


def s10_zero_changes(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='The payoff',
                   title='Zero data-plane changes required')

    rows = [
        ('Indexer service',       'Unchanged',  EMERALD),
        ('Index naming',          'Unchanged',  EMERALD),
        ('Mappings',              'Unchanged',  EMERALD),
        ('Trade-service / Kafka', 'Unchanged',  EMERALD),
        ('OpenSearch UI config',  'New IAM trust + add domain',  AMBER),
        ('Auth (IAM or IDC)',     'Configure once per cluster',  AMBER),
    ]
    y = Inches(2.3); h = Inches(0.55); gap = Inches(0.10)
    for i, (label, change, accent) in enumerate(rows):
        ry = y + (h + gap) * i
        add_rounded_rect(s, Inches(0.6), ry, Inches(8.0), h, fill=WHITE,
                         line=SLATE_200, line_w=Pt(0.5))
        add_rect(s, Inches(0.6), ry, Inches(0.10), h, fill=accent)
        add_textbox(s, Inches(0.95), ry + Inches(0.13), Inches(7.5), Inches(0.4),
                    label, size=14, bold=True, color=SLATE_900)
        add_rounded_rect(s, Inches(8.85), ry, Inches(3.85), h, fill=accent, line=None)
        add_textbox(s, Inches(8.85), ry + Inches(0.13), Inches(3.85), Inches(0.4),
                    change, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_speaker_notes(s,
        "What did we have to change in the project to take advantage of "
        "this? Look at the table. Indexer service — unchanged. Index naming "
        "convention — unchanged. Mappings — unchanged. Trade service and "
        "Kafka pipeline — completely untouched. The only changes are at the "
        "OpenSearch UI configuration layer: a new IAM trust policy on each "
        "regional domain, then add the domains to the UI application. "
        "Configure your auth method once per cluster — IAM or IAM Identity "
        "Center — and you're done. Net code change in this project: zero.")
    return s


def s11_demo_scenario(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Demo scenario',
                   title="A trader in NYC, asking about EU + APAC trades")

    add_rounded_rect(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.0),
                     fill=NAVY_DARK, line=None)
    add_textbox(s, Inches(0.85), Inches(2.35), Inches(11.5), Inches(0.35),
                'THE QUESTION', size=10, bold=True, color=AWS_ORANGE)
    add_textbox(s, Inches(0.85), Inches(2.65), Inches(11.5), Inches(0.45),
        '“Show me all HIGH-risk USD trades, last 24h, across every region.”',
        size=16, color=WHITE, font='Menlo')

    # 3 step layout
    steps = [
        ('Open',     'Single OpenSearch UI bookmark · log in once'),
        ('Query',    'Index pattern: fx-trades-* · risk:HIGH AND from:USD AND ts:[now-24h TO now]'),
        ('Results',  'Hits returned from all 3 regional domains · per-region timing'),
    ]
    y = Inches(3.55); h = Inches(0.7); gap = Inches(0.10)
    for i, (head, body) in enumerate(steps):
        ry = y + (h + gap) * i
        add_rounded_rect(s, Inches(0.6), ry, Inches(2.0), h, fill=AWS_ORANGE, line=None)
        add_textbox(s, Inches(0.6), ry + Inches(0.20), Inches(2.0), Inches(0.4),
                    head.upper(), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rounded_rect(s, Inches(2.75), ry, Inches(9.95), h, fill=WHITE, line=SLATE_200, line_w=Pt(0.5))
        add_textbox(s, Inches(3.00), ry + Inches(0.20), Inches(9.5), Inches(0.4),
                    body, size=13, color=SLATE_900, font='Menlo' if i in (1, 2) else FONT)

    add_textbox(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.4),
        'Three regions queried in parallel. No replicated indices. No bookmarks per region.',
        size=12, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "Concrete demo. Imagine a trader in NYC. The question: show me all "
        "HIGH-risk USD trades from the last 24 hours, anywhere in the world. "
        "Three steps in OpenSearch UI. Open one URL. Log in once. Run one "
        "query against the index pattern fx-trades-star — wildcard. Behind "
        "the scenes the UI fans the query out to all three regional domains "
        "in parallel, returns the merged results, and even shows you the "
        "per-region latency. Same UX as querying a single domain. No "
        "bookmarks, no copy-paste, no replicated indices.")
    return s


def s12_setup(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='How it’s set up',
                   title='Cross-region in three configuration steps')

    steps = [
        ('Pick the home Region',
         'Where the OpenSearch UI application lives. All cross-region calls originate here.'),
        ('Add target domains',
         'In the UI’s data-source config, add each remote domain ARN. Public or VPC, both supported.'),
        ('Grant trust',
         'On every target domain, allow the UI Region’s service principal to read indices.'),
    ]
    y = Inches(2.3); h = Inches(1.05); gap = Inches(0.15)
    for i, (head, body) in enumerate(steps):
        ry = y + (h + gap) * i
        add_rounded_rect(s, Inches(0.6), ry, Inches(0.95), h, fill=AWS_DARK, line=None)
        add_textbox(s, Inches(0.6), ry + Inches(0.35), Inches(0.95), Inches(0.4),
                    f'{i+1}', size=24, bold=True, color=AWS_ORANGE, align=PP_ALIGN.CENTER)
        add_rounded_rect(s, Inches(1.65), ry, Inches(11.05), h, fill=WHITE,
                         line=SLATE_200, line_w=Pt(0.75))
        add_textbox(s, Inches(1.95), ry + Inches(0.18), Inches(10.5), Inches(0.4),
                    head, size=17, bold=True, color=SLATE_900)
        add_textbox(s, Inches(1.95), ry + Inches(0.55), Inches(10.5), Inches(0.5),
                    body, size=13, color=SLATE_500, line_spacing=1.4)

    add_textbox(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
        'Full step-by-step in the AWS docs link on the outro slide.',
        size=12, color=SLATE_500, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "Three configuration steps. One — pick the home Region for the "
        "OpenSearch UI application. All cross-region calls originate from "
        "there, so that's also where your latency is measured from. Two — "
        "in the UI's data source configuration, add each remote OpenSearch "
        "domain by ARN. Both public-internet domains and VPC-attached "
        "domains are supported. Three — on each target domain, grant trust "
        "to the UI's home-Region service principal. That's it. No new "
        "infrastructure, no proxies, no Lambdas in the middle.")
    return s


def s13_iam(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Auth · option 1',
                   title='IAM-based authentication')

    items = [
        'Each user assumes an IAM role mapped to OpenSearch fine-grained access control',
        'Fits orgs that already use IAM federation (SAML / OIDC → IAM roles)',
        'Same authorization logic across regions — IAM trust policies do the routing',
        'Per-index, per-action policies still apply on every regional domain',
    ]
    add_bullets(s, Inches(0.6), Inches(2.2), Inches(12), Inches(4),
                items, size=17, color=SLATE_900, gap_after=Pt(14))

    add_rounded_rect(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.7),
                     fill=NAVY_SOFT, line=None)
    add_textbox(s, Inches(0.95), Inches(6.10), Inches(11.5), Inches(0.5),
        'Best when: your org already runs IAM-managed identities · auditing via CloudTrail',
        size=14, color=NAVY_DARK, bold=True, line_spacing=1.4)
    add_speaker_notes(s,
        "Auth option one — plain IAM. Each user assumes an IAM role that's "
        "mapped to OpenSearch's fine-grained access control. Same model "
        "you've used for years — federation through SAML or OIDC into IAM "
        "roles, then the role gets read access on the indices. Authorization "
        "is consistent across regions because IAM trust policies are the "
        "thing being evaluated. Per-index and per-action restrictions still "
        "apply on every regional domain — cross-region access doesn't bypass "
        "your existing security model. Best fit if your organization already "
        "manages identities through IAM and audits via CloudTrail.")
    return s


def s14_idc(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Auth · option 2',
                   title='IAM Identity Center for SSO')

    items = [
        'Users log in once via IDC — UI federates that to each regional domain',
        'Native group-to-permission mapping — no IAM role choreography',
        'Plays nicely with multi-account orgs (AWS Organizations + IDC)',
        'Audit trail unified through IDC, not stitched across IAM logs',
    ]
    add_bullets(s, Inches(0.6), Inches(2.2), Inches(12), Inches(4),
                items, size=17, color=SLATE_900, gap_after=Pt(14))

    add_rounded_rect(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.7),
                     fill=EMERALD_SOFT, line=None)
    add_textbox(s, Inches(0.95), Inches(6.10), Inches(11.5), Inches(0.5),
        'Best when: you’ve standardized on IDC across the organization · cleaner SSO UX',
        size=14, color=SLATE_900, bold=True, line_spacing=1.4)
    add_speaker_notes(s,
        "Auth option two — IAM Identity Center, formerly AWS SSO. The "
        "user logs in once with their corporate identity. The UI federates "
        "that single sign-on across all the regional domains. Group memberships "
        "map directly to OpenSearch permissions — no need to choreograph IAM "
        "role assumptions for each user. Plays particularly nicely with "
        "AWS Organizations and multi-account setups, which we'll get to on "
        "the next slide. Best fit if you've already standardized on IDC.")
    return s


def s15_xaccount(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Combined power',
                   title='Cross-account + cross-region together')

    # Visual matrix
    add_rounded_rect(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(2.2),
                     fill=AWS_DARK, line=None)
    add_textbox(s, Inches(0.95), Inches(2.4), Inches(11.5), Inches(0.4),
                'EARLIER IN 2026 + NOW = THE FULL ENTERPRISE STORY',
                size=11, bold=True, color=AWS_ORANGE)
    add_textbox(s, Inches(0.95), Inches(2.75), Inches(11.5), Inches(0.6),
        'Single OpenSearch UI · multiple accounts · multiple Regions · any combination.',
        size=20, bold=True, color=WHITE, line_spacing=1.4)
    add_textbox(s, Inches(0.95), Inches(3.5), Inches(11.5), Inches(0.7),
        'Production account · staging account · per-team accounts · per-region replicas — '
        'one pane of glass for all of them.',
        size=14, color=CYAN_SOFT, line_spacing=1.5)

    # 3 personas / use cases
    personas = [
        ('Platform engineer', 'Sees prod + staging + dev in one UI'),
        ('Security analyst',  'Cross-account threat hunting'),
        ('Compliance lead',   'Per-region data still where it should be'),
    ]
    for i, (p, body) in enumerate(personas):
        x = Inches(0.6) + Inches(4.05) * i
        add_rounded_rect(s, x, Inches(4.7), Inches(3.95), Inches(1.7),
                         fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
        add_textbox(s, x + Inches(0.2), Inches(4.85), Inches(3.6), Inches(0.4),
                    p.upper(), size=11, bold=True, color=NAVY_DARK)
        add_textbox(s, x + Inches(0.2), Inches(5.2), Inches(3.6), Inches(1.0),
                    body, size=14, color=SLATE_700, line_spacing=1.5)

    add_speaker_notes(s,
        "And here's where it gets really interesting. Earlier this year, AWS "
        "shipped cross-account data access for OpenSearch UI. Now they've "
        "added cross-region. The combination is the full enterprise story. "
        "A single OpenSearch UI application can talk to OpenSearch domains "
        "in any combination of accounts and regions. Production account, "
        "staging, per-team accounts, per-region replicas — one pane of glass "
        "for all of them. Three personas benefit immediately: platform "
        "engineers see all environments at once; security analysts can do "
        "cross-account threat hunting; compliance leads keep data exactly "
        "where it needs to be while enabling visibility for the people who "
        "need it.")
    return s


def s16_compliance(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Why it matters · compliance',
                   title='Data stays put — visibility doesn’t have to')

    regimes = [
        ('GDPR (EU)',     'EU customer data must stay in EU regions',     CYAN),
        ('MAS (SG)',      'Singapore Monetary Authority — data residency', AMBER),
        ('RBI (IN)',      'India payment-data localization',                EMERALD),
        ('PIPL (CN)',     'China cross-border transfer restrictions',       ROSE),
        ('Industry NDAs', 'Customer contracts pinning data to regions',     INDIGO),
    ]
    y = Inches(2.3); row_h = Inches(0.65); gap = Inches(0.08)
    for i, (regime, body, accent) in enumerate(regimes):
        ry = y + (row_h + gap) * i
        add_rounded_rect(s, Inches(0.6), ry, Inches(2.7), row_h, fill=accent, line=None)
        add_textbox(s, Inches(0.6), ry + Inches(0.18), Inches(2.7), Inches(0.4),
                    regime, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rounded_rect(s, Inches(3.45), ry, Inches(9.25), row_h,
                         fill=WHITE, line=SLATE_200, line_w=Pt(0.5))
        add_textbox(s, Inches(3.7), ry + Inches(0.18), Inches(8.95), Inches(0.4),
                    body, size=13, color=SLATE_900)

    add_textbox(s, Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.4),
        'Cross-region UI federates queries — it does not move data. Compliance unchanged.',
        size=12, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "Why this matters from a compliance angle. GDPR for EU customer "
        "data. The Monetary Authority of Singapore. India's RBI on payment "
        "data localization. China's PIPL on cross-border transfers. Plus "
        "the industry-specific NDAs and customer contracts that pin data to "
        "regions. Every one of these regimes says: data must stay in its "
        "origin region. The critical thing about cross-region OpenSearch UI "
        "is that it federates queries — it doesn't move data. The bytes "
        "stay where they were written. Only query results travel back to "
        "the user, which is exactly the compliant behavior these regimes "
        "are looking for.")
    return s


def s17_cost(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Why it matters · economics',
                   title='Inter-region egress · the silent budget killer')

    # Compare two columns
    add_rounded_rect(s, Inches(0.6), Inches(2.2), Inches(5.95), Inches(4.4),
                     fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
    add_rect(s, Inches(0.6), Inches(2.2), Inches(0.10), Inches(4.4), fill=ROSE)
    add_textbox(s, Inches(0.95), Inches(2.4), Inches(5.5), Inches(0.4),
                'BEFORE · CROSS-CLUSTER REPLICATION', size=11, bold=True, color=ROSE)
    add_textbox(s, Inches(0.95), Inches(2.75), Inches(5.5), Inches(0.5),
                'You pay for every byte that crosses regions',
                size=16, bold=True, color=SLATE_900)
    add_bullets(s, Inches(0.95), Inches(3.45), Inches(5.4), Inches(3.0),
                ['~$0.02 / GB inter-region egress',
                 'High-write workloads = thousands per month',
                 'Plus extra storage on the central region',
                 'Plus replication lag to manage'],
                size=13, color=SLATE_700, gap_after=Pt(8), bullet_color=ROSE)

    add_rounded_rect(s, Inches(6.85), Inches(2.2), Inches(5.95), Inches(4.4),
                     fill=WHITE, line=SLATE_200, line_w=Pt(0.75))
    add_rect(s, Inches(6.85), Inches(2.2), Inches(0.10), Inches(4.4), fill=EMERALD)
    add_textbox(s, Inches(7.20), Inches(2.4), Inches(5.5), Inches(0.4),
                'AFTER · CROSS-REGION UI', size=11, bold=True, color=EMERALD)
    add_textbox(s, Inches(7.20), Inches(2.75), Inches(5.5), Inches(0.5),
                'You only pay for query result bytes',
                size=16, bold=True, color=SLATE_900)
    add_bullets(s, Inches(7.20), Inches(3.45), Inches(5.4), Inches(3.0),
                ['Storage stays in origin Region',
                 'Egress = small result payloads only',
                 'No replication infra to manage',
                 'Cost scales with USE, not data size'],
                size=13, color=SLATE_700, gap_after=Pt(8), bullet_color=EMERALD)

    add_speaker_notes(s,
        "The economics. Inter-region data transfer is one of those AWS line "
        "items that creeps up on you. Roughly two cents per gigabyte going "
        "between most regions. That sounds small until you have a "
        "high-throughput workload constantly replicating to a central "
        "region — then it adds up to thousands per month, plus the doubled "
        "storage cost, plus you have to manage replication lag. With "
        "cross-region OpenSearch UI, storage stays in the origin region — "
        "zero replication egress. The only bytes that cross regions are "
        "query result payloads, which are tiny. Cost scales with how often "
        "people query, not with how much data you have. For most workloads "
        "that's a one-to-two order-of-magnitude reduction.")
    return s


def s18_try_it(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow='Try it yourself',
                   title='Get hands-on in 30 minutes')

    steps = [
        '1. Spin up the FX project locally — one npm command (link in description)',
        '2. Deploy the included GitHub Actions workflows to your AWS account',
        '3. Replace local OpenSearch with managed in 2 regions',
        '4. Wire OpenSearch UI to both domains following the AWS doc link',
        '5. Run the demo query — fx-trades-* — and watch the federation happen',
    ]
    add_bullets(s, Inches(0.6), Inches(2.3), Inches(12), Inches(4.4),
                steps, size=18, color=SLATE_900, gap_after=Pt(16),
                bullet_color=AWS_ORANGE)

    add_rounded_rect(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.7),
                     fill=AWS_DARK, line=None)
    add_textbox(s, Inches(0.95), Inches(6.20), Inches(11.5), Inches(0.4),
        'Repo + AWS doc link on the outro slide.',
        size=14, color=AWS_ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "If you want to try this yourself, here's the path. Five steps. "
        "Spin up the FX project locally — one npm command, fully containerized. "
        "Deploy the included GitHub Actions workflows to your own AWS "
        "account; the VPC and IAM scaffolding is already there. Swap local "
        "OpenSearch for managed OpenSearch domains in two regions of your "
        "choice. Wire OpenSearch UI to both domains following the AWS docs "
        "link. Then run the demo query — fx-trades-star — and watch the "
        "federation happen. Should take a developer about thirty minutes "
        "end-to-end if you already have an AWS account.")
    return s


def s19_roadmap(prs, page, total):
    s = base_slide(prs, page=page, total=total,
                   eyebrow="What's next on the channel",
                   title='Coming up')

    eps = [
        ('Next', 'Walking through the OpenSearch UI cross-region setup live',  AWS_ORANGE),
        ('Then', 'Deploying the FX project to AWS — full GitHub Actions tour', NAVY),
        ('Then', 'Cross-account demo — multi-account organization setup',      CYAN),
        ('Then', 'Querying across regions for compliance + observability',     EMERALD),
    ]
    y = Inches(2.4); h = Inches(0.85); gap = Inches(0.12)
    for i, (label, body, accent) in enumerate(eps):
        ry = y + (h + gap) * i
        add_rounded_rect(s, Inches(0.6), ry, Inches(12.1), h, fill=WHITE,
                         line=SLATE_200, line_w=Pt(0.5))
        add_rounded_rect(s, Inches(0.85), ry + Inches(0.20), Inches(1.6), Inches(0.45),
                         fill=accent, line=None)
        add_textbox(s, Inches(0.85), ry + Inches(0.24), Inches(1.6), Inches(0.4),
                    label.upper(), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(2.65), ry + Inches(0.25), Inches(9.8), Inches(0.45),
                    body, size=15, color=SLATE_900)
    add_speaker_notes(s,
        "Coming up on the channel. Next video — I do the OpenSearch UI "
        "cross-region setup live, on real domains. Then we deploy the FX "
        "project end-to-end to AWS, walking through the GitHub Actions "
        "workflows. Then a cross-account demo using AWS Organizations. "
        "And finally a video on the actual queries that solve real "
        "compliance and observability problems. Subscribe so you don't "
        "miss the live setup video — that's where the rubber meets the road.")
    return s


def s20_outro(prs, _total):
    s = base_slide(prs, no_chrome=True)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=AWS_DARK)
    add_rect(s, 0, 0, Inches(7.0), SLIDE_H, fill=NAVY_DARK)
    add_rect(s, 0, Inches(7.30), SLIDE_W, Inches(0.20), fill=AWS_ORANGE)
    add_rect(s, 0, Inches(7.20), Inches(5.0), Inches(0.10), fill=CYAN)

    add_textbox(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                'THANKS FOR WATCHING', size=14, bold=True, color=AWS_ORANGE)
    add_textbox(s, Inches(0.8), Inches(2.0), Inches(11), Inches(1.4),
                'Like · Subscribe · Comment', size=46, bold=True, color=WHITE)
    add_textbox(s, Inches(0.8), Inches(3.5), Inches(11), Inches(0.6),
                'Tell me what you want me to build live next — drop it in the comments.',
                size=18, color=CYAN_SOFT)

    # Two link cards
    add_rounded_rect(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.85),
                     fill=NAVY_DARK, line=AWS_ORANGE, line_w=Pt(1.0))
    add_textbox(s, Inches(1.0), Inches(4.60), Inches(11.0), Inches(0.4),
                'AWS WHAT’S NEW · OFFICIAL ANNOUNCEMENT', size=10, bold=True, color=AWS_ORANGE)
    add_textbox(s, Inches(1.0), Inches(4.85), Inches(11.0), Inches(0.5),
        'aws.amazon.com/about-aws/whats-new/2026/05/opensearch-ui-cross-region-data-access-domains',
        size=12, color=WHITE, font='Menlo')

    add_rounded_rect(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.85),
                     fill=NAVY_DARK, line=CYAN, line_w=Pt(1.0))
    add_textbox(s, Inches(1.0), Inches(5.60), Inches(11.0), Inches(0.4),
                'WORKED EXAMPLE · SOURCE CODE', size=10, bold=True, color=CYAN)
    add_textbox(s, Inches(1.0), Inches(5.85), Inches(11.0), Inches(0.5),
                'github.com/javakishore-veleti/FX-Trade-Analytics-AWs-OpenSearch',
                size=12, color=WHITE, font='Menlo')

    add_textbox(s, Inches(0.8), Inches(6.85), Inches(11), Inches(0.4),
                "Channel: Training My Hobby · Series: AWS What's New",
                size=11, color=SLATE_400)
    add_speaker_notes(s,
        "If this walkthrough helped, hit like and subscribe — it genuinely "
        "moves the channel. Two links on screen. The official AWS what's-new "
        "post for the announcement. The source code repo for the FX worked "
        "example — go fork it. Drop me a comment with what you want me to "
        "build live next: the cross-region setup itself, the AWS deploy, "
        "the cross-account demo, or something else entirely. Thanks for "
        "watching, see you in the next one.")
    return s


# ---------------------------------------------------------------- main
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        (s01_title,           False),
        (s02_announcement,    True),
        (s03_problem,         True),
        (s04_before,          True),
        (s05_after,           True),
        (s06_architecture,    True),   # merged: project intro + topology
        (s08_partitioned,     True),
        (s09_forward_looking, True),
        (s10_zero_changes,    True),
        (s11_demo_scenario,   True),
        (s12_setup,           True),
        (s13_iam,             True),
        (s14_idc,             True),
        (s15_xaccount,        True),
        (s16_compliance,      True),
        (s17_cost,            True),
        (s18_try_it,          True),
        (s19_roadmap,         True),
        (s20_outro,           False),
    ]
    total = len(builders)
    for i, (build, takes_page) in enumerate(builders, 1):
        if takes_page:
            build(prs, i, total)
        else:
            build(prs, total)

    out = Path(__file__).parent / 'AWS-OpenSearch-UI-Cross-Region-FX-Demo.pptx'
    prs.save(out)
    print(f'Wrote {out}')
    print(f'Slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()
